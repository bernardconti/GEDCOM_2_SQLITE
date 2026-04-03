import os
import os.path
import sys
sys.path.append(os.path.dirname(__file__))
from MH_MODULE_SQLQUERIES_1 import *
from MH_MODULE_UTILS_1 import *
from MH_MODULE_PPTX_1 import *

# file management

from io import BytesIO
from datetime import date,datetime
from pytz import timezone as tz

import matplotlib.pyplot as plt

#-- _PPTX_
#pip install pypptx
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.shapes import *
from pptx.enum.text import PP_ALIGN,MSO_AUTO_SIZE
from pptx.dml.color import *
from pptx.enum.dml import *
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.table import _Cell, _Row

#pip install pypptx-with-oxml
from pptx.oxml import parse_xml  ## for table backgroud
from pptx.oxml.ns import nsdecls  ## for table bakcgroud

from pptx.dml.color import RGBColor
#import textwrap
import copy
#from collections import Counter
from io import BytesIO
# --------------------------------------------------------------
from html.parser import HTMLParser
# --------------------------------------------------------------
icloud = "/Users/bernardconti/Library/Mobile Documents/com~apple~CloudDocs"
les_WIP = "/Users/bernardconti/LOCAL_TEMP/WIP/"
les_PDF_RV = "/Users/bernardconti/Library/Mobile Documents/com~apple~CloudDocs/Les éditions du 57/PDF_RectoVerso/"
icon_dir = icloud+'/MesProgrammes/python_global/icons'
export_path =  "/Users/bernardconti/LOCAL_TEMP/Photos"
save_path = "/Users/bernardconti/LOCAL_TEMP/Documents"
import_dir = "/Users/bernardconti/LOCAL_TEMP/ePhoto_Doc/Import/"
watermarkdir = "/Users/bernardconti/LOCAL_TEMP/watermark"
le_MH_Photos_Bio = icloud+'/MesProgrammes/MH_Photos_Bio/'
#--------------------------------------------------------------
retour_ligne = '\n'
Gray= "#EFEEEF"
GrayRow= "#EFEEEF"
Black ="#000000"
White = "#FFFFFF"
couleur_homme ="#ffcc99"
couleur_femme="#ccccff"
font_texte ="arial narrow"
couleur_cible = "#30DD30"
Gray1= "#b7b1b1"#F2F2F2
GraySide = "#c5c5c5"
Gray3 ="#D9D6D6"
Gray4 = "#E1DEDEEF"
couleur_chemin = "#EC5800"
couleur_titre ="#EB742F"
Green = "#B4E1C0"
Gray2= "#fde9d9"
Silver = "#A9A6A6" 
cell11 = "#F8CEB8"
cell12 = "#D9D9D9"
cell21 = "#FDE8DD"
cell22 = "#F2F2F2"
#Bleu = "#375e94"
Bleu = "#272170"
Darkblue = "#166082"
Gris_clair = "#EFEFEF"
ligne = "#375e94"
descendant = "#000000"
red = "#FF0000"
page_width = 190
pt_mm = 0.68
document_font = "Aptos (Corps)"
font_file = "/Users/bernardconti/Library/Fonts/Aptos.ttf"
font_file_bold = "/Users/bernardconti/Library/Fonts/Aptos-Bold.ttf"

EMU = 36000

#PPT
slide_width = 277
slide_height = 210
slide_margin_left = 5
slide_margin_right = 5
slide_margin_top = 10
slide_margin_bottom = 5
image_size = 28.2

slide_layout_box_0 = 0
slide_layout_box_0_1 = 1
slide_layout_ascendants_4 = 2
slide_layout_ascendants_tous = 3
slide_layout_couples_liste = 4
slide_layout_descendants_draw = 4
slide_layout_liste = 5
slide_layout_57 = 6
slide_layout_garde = 7
slide_layout_photoMH = 8
slide_layout_album = 9
slide_layout_entourage = 10
slide_layout_descendants_list = 11
slide_layout_section = 12
slide_layout_descendants_details = 13
slide_layout_table_image = 14
slide_layout_table_pleine = 15
#photos
box_photo_width = 50
box_photo_height = 50


#### PPTX SECTION
#=============================================================================================================
# PAGES
#=============================================================================================================
def hex_to_rgb(value):
#-------------------------------------------------------------------------------------------------------------
    RGB = RGBColor(0, 0, 0)
    try:
        value = value.lstrip('#')
        lv = len(value)
        (rouge, vert, bleu) = tuple(int(value[i:i + lv // 3], 16) for i in range(0, lv, lv // 3))
        RGB = RGBColor(rouge, vert, bleu)
    except Exception as error:
        print (value,tuple(int(value[i:i + lv // 3], 16) for i in range(0, lv, lv // 3)))
#-------------------------------------------------------------------------------------------------------------
    return RGB
#============================================================================================================= 
def PPTX_add_page(le_document,le_layout):
    boxes = []
    la_page = le_document.slides.add_slide(le_document.slide_layouts[le_layout])
    for shape in la_page.placeholders:
        #print('%d %s' % (shape.placeholder_format.idx, shape.name))
        boxes.append(shape)
    b = PPTX_add_box(la_page,slide_width-10,slide_height-5,10,5)
    n_slide = le_document.slides.index(la_page)
    if n_slide > 1 : PPTX_add_run(PPTX_add_paragraph(b,"right"),str(n_slide),"italic","bold",size=10)
    return la_page,boxes
#=============================================================================================================
def PPTX_add_page_garde(sql_obj,le_document,les_bros,isModeFratrie):
    la_page,boxes = PPTX_add_page( le_document,slide_layout_garde)
    p = PPTX_add_paragraph(boxes[0])

    boxes[0].text = f'{datetime.now().month}/{datetime.now().year}'
    if isModeFratrie:
        for idx,le_bro in enumerate(les_bros):
            p = PPTX_add_paragraph(boxes[0])
            PPTX_add_run(p,text_personne_full(le_bro))
            PPTX_add_photoID(sql_obj,la_page,20 + 40*idx,20,40,40,le_bro)
    else: 
        p = PPTX_add_paragraph(boxes[0])
        PPTX_add_run(p,text_personne_full(les_bros[0]))
        PPTX_add_photoID(sql_obj,la_page,20,20,40,40,les_bros[0])

    return la_page,boxes
#=============================================================================================================
def PPTX_page_image(sql_obj,le_document,MH_personne,sujet,*args,**kwargs):
#=============================================================================================================
    isPhotoID =  True
    le_layout = slide_layout_entourage
    x0 = slide_margin_left
    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
            if valeur == "nophotoid" : isPhotoID = False 
            if valeur == "photomh" : le_layout = slide_layout_photoMH
            if valeur == "liste": le_layout = slide_layout_liste
            if valeur == "image_right" : x0 = slide_width  - image_size

    le_layout = slide_layout_entourage
    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "layout" : le_layout = valeur
    # ajout d'une slide

    la_page, boxes = PPTX_add_page(le_document,le_layout)
    if isPhotoID : PPTX_add_photoID(sql_obj,la_page,x0,0,image_size,image_size,MH_personne)
    # Nom Premon  
    p = PPTX_add_paragraph(boxes[0],"center")
    PPTX_add_run(p,f'{MH_personne.prenom} {MH_personne.nom}',"bold")
    PPTX_add_run(p,f' {MH_personne.prenoms if MH_personne.prenoms else ""}',"italic",)
    PPTX_add_run(p,f' "{MH_personne.surnom if MH_personne.surnom else ""}"',"italic")

    if sujet:
        p = PPTX_add_paragraph(boxes[-1],"center")
        PPTX_add_run(p,sujet)
    return la_page, boxes
#=============================================================================================================
def PPTX_add_page_section(le_document,MH_adult1,MH_adult2,le_sujet,box_sommaire):

    if box_sommaire and le_sujet:
        la_page,boxes = PPTX_add_page( le_document,slide_layout_garde)

        if isinstance(le_sujet, list): le_texte = [" ".join(le_sujet)]
        elif isinstance(le_sujet, str): le_texte = [le_sujet]
        else: return

        if MH_adult1 != MH_none:le_texte.append(f'de {MH_adult1.prenom} {MH_adult1.nom}')
        if MH_adult2 != MH_none:le_texte.append(f' et de {MH_adult2.prenom} {MH_adult2.nom}')
        for item in le_texte:
            PPTX_add_run(PPTX_add_paragraph(boxes[0]),item)

        PPTX_add_ligne_sommaire(box_sommaire,le_document,la_page,le_texte,"section")

    return 
#=============================================================================================================
# TABLES
#=============================================================================================================
def PPTX_table_add(la_page,col_format,x,y,w,h):
#------------------------------------------------------------------------------------------------------------- 
    la_table = None
    if w >0 and h > 0 and la_page and col_format:  

        if len(col_format) > 0 : 
            n_row = 1
    # ajout table
            la_table  = la_page.shapes.add_table(n_row,len(col_format),Mm(x),Mm(y),Mm(w),Mm(h)).table
            la_table.first_row = False
            for idx,item in enumerate(col_format):
                if item[0] : la_table.columns[idx].width = Mm(item[0])
    else:
        print("PPTX_table_add",w,h,la_page,col_format)

    return la_table
#=============================================================================================================
def PPTX_table_display(sql_obj,le_document,col_format,le_header,display_table,box_sommaire,*args,**kwargs):
#-------------------------------------------------------------------------------------------------------------
    x_table = 10
    y_table = 0
    le_layout = slide_layout_table_image
    for clef, valeur in kwargs.items():  
        #if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "x_table" : x_table = valeur
            if clef == "y_table" : y_table = valeur
#-------------------------------------------------------------------------------------------------------------          
#   display_table   [
#   table_row           [ 
#                           isNew,                          table_row[0]
#                           isEven,                         table_row[1]
#                           max_n_col,                      table_row[2]
#   table_cells             [                               table_row[3]
#   table_cell                 col_format[0] + [texte],     texte = string or list
#                              col_format[1] + [texte], 
#                              ...
#                           ],
#                           MH_persomme pour photoID        table_row[4]
#                       ]   
#                   ]  
#-------------------------------------------------------------------------------------------------------------
#   col_format     [    0:width,1:justif,2:font_type,3:font_size:
#                       4:text color even,5:back ground color even
#                       6:text color odd,7:back ground color odd
#                       8:margin_left,9:margin_right,10:margin_top,11:margin_bottom 
#                   ]
#-------------------------------------------------------------------------------------------------------------
    htable = 0
    h_row = 0
    la_table = None
    les_PhotoIDs = []
    #---------------------------------------------------------------------------------------------------------
    for table_row in display_table:
    #---------------------------------------------------------------------------------------------------------
        isNew = table_row[0]
        # compute table height after new row 
        h_row = PPTX_table_row_height(table_row[3])
        htable = htable + h_row

        if htable > slide_height or isNew:
            if les_PhotoIDs :
                PPTX_add_photoIDs_side(sql_obj,la_page,reversed(les_PhotoIDs),"image",box_width = x_table)
                les_PhotoIDs = []

            la_page, boxes = PPTX_add_page(le_document,le_layout) 
            PPTX_add_box(la_page,0,0,x_table,slide_height,bcolor = GraySide)

            la_table = PPTX_table_add(la_page,col_format,x_table,y_table,1,1)
            if le_header : 
                PPTX_table_row_display_photoID_in_cell(sql_obj,la_table,le_header,"header")  
                htable = PPTX_table_row_height(le_header[3]) + h_row

        # PhotoID
        if table_row[5]: les_PhotoIDs.append(table_row[5])
#----------------------------------------------------------------------------------------------------------    
        PPTX_table_row_display_photoID_in_cell(sql_obj,la_table,table_row)

    if les_PhotoIDs: PPTX_add_photoIDs_side(sql_obj,la_page,reversed(les_PhotoIDs),"image",box_width = x_table)

    return
#=============================================================================================================
def PPTX_table_display_photoID_in_cell(sql_obj,le_document,col_format,le_header,display_table,box_sommaire,le_titre,*args,**kwargs):
#-------------------------------------------------------------------------------------------------------------
    x_table = 0
    y_table = 0
    le_layout = slide_layout_table_image    
    for clef, valeur in kwargs.items():  
        #if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "x_table" : x_table = valeur
            if clef == "y_table" : y_table = valeur
#-------------------------------------------------------------------------------------------------------------          
    htable = 0
    h_row = 0
    la_table = None
    isFirst = True
#-------------------------------------------------------------------------------------------------------------
    for table_row in display_table:
#-------------------------------------------------------------------------------------------------------------
        isNew = table_row[0]
        # compute table height after new row 
        #if table_row[3][0][12] == "photoID": 
        h_row = PPTX_table_row_height(table_row[3])
        htable = htable + h_row
#---------------------------------------------------------------------------------------------------------- 
        if htable > slide_height or isNew:
           
            la_page, boxes = PPTX_add_page(le_document,le_layout) 
            PPTX_add_box(la_page,0,0,x_table,slide_height,bcolor = GraySide)

            if box_sommaire and isFirst and le_titre: 
                PPTX_add_ligne_sommaire(box_sommaire,le_document,la_page,le_titre)
                isFirst = False

            la_table = PPTX_table_add(la_page,col_format,x_table,y_table,1,1)
            if le_header : 
                PPTX_table_row_display_photoID_in_cell(sql_obj,la_table,le_header,"header")  
                htable = PPTX_table_row_height(le_header[3]) + h_row

#----------------------------------------------------------------------------------------------------------    
        PPTX_table_row_display_photoID_in_cell(sql_obj,la_table,table_row,*args)
    return
#=============================================================================================================
def PPTX_table_row_add(table):
#----------------------------------------------------------------------------------------------------------
    new_row = copy.deepcopy(table._tbl.tr_lst[0])  # copies last row element

    for tc in new_row.tc_lst:
        cell = _Cell(tc, new_row.tc_lst)
        cell.text = ''

    table._tbl.append(new_row)
    row = _Row(new_row, table)

    for tc in row.cells:
        if tc.is_merge_origin : tc.split()

    return row,row._tr.row_idx
#=============================================================================================================
def PPTX_table_row_display_photoID_in_cell(sql_obj,la_table,table_row,*args):
#-----------------------------------------------------------------------------------------------------------  
    isEven = table_row[1]
    max_n_col = table_row[2]
    table_cells = table_row[3]
    n_col = len(table_cells)-1
    n_FixLast = table_row[4]
    #--------------------------------------------------------------------------------------------------------- 
    isHeader = False  
    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
            if valeur == "header" : isHeader =  True
    # a new row to the table
    if isHeader : 
        n_last = 0
    else:
        new_row,n_last = PPTX_table_row_add(la_table)
#----------------------------------------------------------------------------------------------------------   
    la_h = 0.1
    for table_cell in table_cells: 
        le_content = table_cell[12]
        if le_content :
            if not isinstance(le_content,str) and not isinstance(le_content,list) and not isinstance(le_content,int):la_h = table_cell[0]
    la_table.rows[n_last].height = Mm(la_h)

    #la_table.rows[n_last].height = Mm(PPTX_table_row_height(table_cells))
#----------------------------------------------------------------------------------------------------------  
    c1 = n_col - n_FixLast
    c2 = max_n_col - n_FixLast
    if c1 < c2: la_table.cell(n_last,c1).merge(la_table.cell(n_last, c2))
#----------------------------------------------------------------------------------------------------------
    for idx,table_cell in enumerate(table_cells):
        if idx > c1 : n_cell = idx + (c2-c1)
        else : n_cell  = idx
        la_cell = la_table.cell(n_last,n_cell)
        la_cell.fill.solid()
        if isEven : la_cell.fill.fore_color.rgb = hex_to_rgb(table_cell[5])
        else : la_cell.fill.fore_color.rgb = hex_to_rgb(table_cell[7])
        la_cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        le_content = table_cell[12]
        if le_content:
            if not isinstance(le_content,str) and not isinstance(le_content,list) and not isinstance(le_content,int): 
                #le_content = MH individual
                try: PPTX_add_table_add_image_in_cell(sql_obj,la_cell,le_content)
                except Exception as error : print("Add Photo in Table Cell",error)
            else:
                les_textes = []
                if isinstance(le_content,str):       les_textes = [le_content]
                elif isinstance(le_content,int):    les_textes = [str(le_content)]
                elif isinstance(le_content,list):    les_textes = le_content

                if les_textes: 

                    p = PPTX_table_cell_add_paragraph(la_cell,table_cell[1],
                        margin_left = table_cell[8],
                        margin_right = table_cell[9],
                        margin_top = table_cell[10],    
                        margin_bottom = table_cell[11]
                        )

                    for idx,le_texte in enumerate(les_textes):
                        if le_texte:
                            if idx > 0 : PPTX_add_run(p,"\n")
                            PPTX_add_run_html(p,f'{le_texte}',
                                table_cell[2], # bold,normal,,...
                                size = table_cell[3],
                                color = table_cell[4] if not isEven else table_cell[6])

                            #PPTX_add_run(p,f'{le_texte}',
                            #    table_cell[2], # bold,normal,,...
                            #    size = table_cell[3],
                            #    color = table_cell[4] if not isEven else table_cell[6])
#-------------------------------------------------------------------------------------------------------------
    return 
#=============================================================================================================
def PPTX_table_row_height(table_cells):
#row_text  [   0:width, 
#               1:justif,
#               2:font_type,
#               3:font_size:
#               4:text color even,
#               5:back ground color even
#               6:text color even,
#               7:back ground color even
#               8:margin_left, 9:margin_right, 10:margin_top, 11:margin_bottom 
#               12:texte string or list
#           ]

    h_row = 0
    # gestion des tables
    
    for table_cell in table_cells :

        la_cell_width = table_cell[0] -  table_cell[8] - table_cell[9]

        if table_cell[2] == "bold" : la_font = ImageFont.truetype(font_file_bold,table_cell[3])
        else:la_font = ImageFont.truetype(font_file,table_cell[3])

        ascent, descent = la_font.getmetrics()

        le_content = table_cell[12]
        if le_content:
            le_h = 0

            if not isinstance(le_content,str) and not isinstance(le_content,int) and not isinstance(le_content,list):
                le_h = table_cell[0]
            else:
                les_textes = []
                if isinstance(table_cell[12],str):       les_textes = [le_content]
                elif isinstance(table_cell[12],list):    les_textes = le_content
                elif isinstance(table_cell[12],int):    les_textes = [str(le_content)]

                if les_textes : 
                    le_h = 0
                    for le_texte in les_textes:
                        les_textes_sp= le_texte.splitlines()
                        for item in les_textes_sp:
                            w = la_font.getlength(item) * 0.352778
                            n = int(w/(la_cell_width))+1
                            le_h = (ascent + descent)*0.352778*1*n
                            #le_h = le_h + h_cell

                    le_h = le_h + table_cell[10] + table_cell[11]

            h_row = max(h_row,le_h  ) 
    #print(h_row)
    return h_row
#=============================================================================================================
def PPTX_table_cell_add_paragraph(la_cell,*args,**kwargs):
#------------------------------------------------------------------------------------------------------------- 
    align = PP_ALIGN.LEFT
    la_margin_bottom = 1
    la_margin_top = 1
    la_margin_left = 1
    la_margin_right = 1
    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
            if valeur == "center" : align =  PP_ALIGN.CENTER    
            if valeur == "right" : align =  PP_ALIGN.RIGHT

    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "bcolor" : la_couleur = valeur
            if clef == "margin_bottom" : la_margin_bottom = valeur
            if clef == "margin_top" : la_margin_top = valeur
            if clef == "margin_left" : la_margin_left = valeur
            if clef == "margin_right" : la_margin_right = valeur 
#-------------------------------------------------------------------------------------------------------------
    la_cell.margin_bottom = Mm(la_margin_bottom)
    la_cell.margin_top = Mm(la_margin_top)
    la_cell.margin_left = Mm(la_margin_left)
    la_cell.margin_right = Mm(la_margin_right)
    text_frame = la_cell.text_frame
    text_frame.clear()  # Clear existing text
    if len(text_frame.paragraphs[0].text) == 0 : le_paragraph = text_frame.paragraphs[0] 
    else: le_paragraph = text_frame.add_paragraph()
    le_paragraph.alignment = align
#-------------------------------------------------------------------------------------------------------------   
    return le_paragraph
#=============================================================================================================
def PPTX_add_table_add_image_in_cell(sql_obj,cell, MH_personne):

    img_path = get_personne_photoID_file(sql_obj,MH_personne)
    img = open(img_path,'rb').read()
    img_bytes = BytesIO(img)

    image_part, rId = cell.part.get_or_add_image_part(img_bytes)
    tcPr = cell._tc.get_or_add_tcPr()

    xml = f"""<a:blipFill {nsdecls("a")}> 
                      <a:blip {nsdecls("r")} r:embed="{rId}"/>
                      <a:stretch>
                        <a:fillRect/>
                      </a:stretch>
                    </a:blipFill>"""
    ele = parse_xml(xml)

    tcPr.clear()
    tcPr.append(ele)
#=============================================================================================================
# BOXES
#=============================================================================================================
def PPTX_add_box(la_page,x,y,w,h,*args,**kwargs):
#------------------------------------------------------------------------------------------------------------- 
    if w < 0 or h < 0 or not la_page :  
        print("PPTX_add_box",w,h)
        return False
#
# lecture des parametres
#
    la_couleur = False
    le_alignement = MSO_ANCHOR.MIDDLE
    la_margin_bottom = 0.5
    la_margin_top = 0.5
    la_margin_left = 0.5
    la_margin_right = 0.5
        
    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "bcolor" : la_couleur = valeur
            if clef == "margin_bottom" : la_margin_bottom = valeur
            if clef == "margin_top" : la_margin_top = valeur
            if clef == "margin_left" : la_margin_left = valeur
            if clef == "margin_right" : la_margin_right = valeur

    isTransparent = False
    isWrap = True
    isAutosize = False
    isCadre = False

    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
            if valeur == "nowrap" : isWrap =  False    
            if valeur == "transparent" : isTransparent =  True  
            if valeur == "autosize" : isAutosize =  True    
            if valeur == "cadre" : isCadre =  True 
#
# ajout box
#
    txBox = la_page.shapes.add_textbox(Mm(x),Mm(y),Mm(w),Mm(h))

    if isCadre : 
        line = txBox.line
        line.color.rgb = RGBColor(200,200,200)
    elif la_couleur:
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = hex_to_rgb(la_couleur)

    if isTransparent:
                            
        def SubElement(parent, tagname, **kwargs):
            element = OxmlElement(tagname)
            element.attrib.update(kwargs)
            parent.append(element)
            return element

        def _set_shape_transparency(shape, alpha):
        # Set the transparency (alpha) of a shape
            ts = shape.fill._xPr.solidFill
            sF = ts.get_or_change_to_srgbClr()
            sE = SubElement(sF, 'a:alpha', val=str(alpha))

        _set_shape_transparency(txBox,44000)

    tf = txBox.text_frame
    tf.word_wrap = isWrap
    tf.auto_size = isAutosize 
    if le_alignement : tf.vertical_anchor = le_alignement
    tf.margin_bottom = Mm(la_margin_bottom)
    tf.margin_top = Mm(la_margin_top)
    tf.margin_left = Mm(la_margin_left)
    tf.margin_right = Mm(la_margin_right)
#-------------------------------------------------------------------------------------------------------------
    return txBox
#=============================================================================================================
def PPTX_add_paragraph(text_Box,*args,**kwargs):
#------------------------------------------------------------------------------------------------------------- 
    align = PP_ALIGN.LEFT
    level = 0
    before_space = 0

    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
            if valeur == "center" : align =  PP_ALIGN.CENTER    
            if valeur == "right" : align =  PP_ALIGN.RIGHT           

    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(valeur,str): clef = clef.lower()
            if clef == "level" : level = int(valeur)
            if clef == "before_space" : before_space = int(valeur)
#-------------------------------------------------------------------------------------------------------------  
    tf = text_Box.text_frame
    if len(tf.paragraphs[0].text) == 0 : le_paragraph = tf.paragraphs[0] 
    else: le_paragraph = tf.add_paragraph()

    le_paragraph.level = level 
    le_paragraph.alignment = align
    le_paragraph.space_before = Pt(before_space)

    return le_paragraph
#=============================================================================================================
def PPTX_add_run(le_paragraph,le_texte,*args,**kwargs):
#------------------------------------------------------------------------------------------------------------- 
    run = None
    if le_texte:
#------------------------------------------------------------------------------------------------------------- 
        box_width = False
        isItalic =False
        isBold = False
        isUnderline = False
        font = False
        size = False
        color = False
        
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "box_width" : box_width = valeur
                if clef == "size" : size = int(valeur)
                if clef == "font" : font = valeur
                if clef == "color" : color = valeur

        for valeur in args:
            if valeur:
                if isinstance(valeur, str): valeur =valeur.lower()
                if valeur == "bold" : isBold =  True    
                if valeur == "italic" : isItalic =  True  
                if valeur == "underline" : isUnderline =  True    
#-------------------------------------------------------------------------------------------------------------   
        if box_width:
            w_chr = (box_width-1)/len(le_texte)
            size = w_chr * 5
            if size > 12.5:size = 12.5
            if size < 8: size = 8
#------------------------------------------------------------------------------------------------------------- 
        run = le_paragraph.add_run()
        run.text =le_texte
        if font : run.font.name = font
        if size : run.font.size = Pt(size)
        if color : run.font.color.rgb = hex_to_rgb(color)        
        if isBold : run.font.bold = isBold
        if isItalic : run.font.italic = isItalic
        if isUnderline : run.font.underline = isUnderline

    return run
#=============================================================================================================
def PPTX_add_hyperlink_to_slide(box, run , n_slide): #15/01/2026
    box.click_action.target_slide = n_slide
    run.hyperlink.address = box.click_action.hyperlink.address
    run.hyperlink._hlinkClick.action = box.click_action.hyperlink._hlink.action
    run.hyperlink._hlinkClick.rId = box.click_action.hyperlink._hlink.rId
    box.click_action.target_slide = None
    return
#=============================================================================================================
def box_font_size(wb,le_texte,size_max,size_min):
    size = size_max
    if wb > 0 and le_texte:
        w_chr = (wb-1)/len(le_texte)
        size = w_chr * 5
        if size > size_max:size = size_max
        if size < size_min: size = size_min
    return size
#=============================================================================================================
def PPTX_add_run_MH_personne(le_paragraph,MH_personne,*args,**kwargs):
#-------------------------------------------------------------------------------------------------------------   
    if MH_personne == MH_none: return
#-------------------------------------------------------------------------------------------------------------  
    color = False
    size = False
    font= False
    prefix = False
    suffix= False

    isUnderline = False

    isNom = False
    isPrenom = False
    isPrenoms = False
    isSurnom = False
    isDate = False
    type_date = 0
    isYear = False
    isbLocation = False
    isdLocation = False
    isCause = False

    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
# lié au run
            if valeur == "underline" : isUnderline =  True
# lié au contenu affiché
            if valeur == "nom" : isNom =  True
            if valeur == "prenom" : isPrenom =  True
            if valeur == "prenoms" : isPrenoms =  True
            if valeur == "surnom" : isSurnom =  True
            if valeur == "date" : isDate =  True
            if valeur == "year" : isYear =  True
            if valeur == "blocation" : isbLocation =  True
            if valeur == "dlocation" : isdLocation =  True
            if valeur == "cause" : isCause =  True

# lié au run
    for clef, valeur in kwargs.items():   
        if clef: 
            clef =clef.lower()
            if clef == "color" : color = valeur
            if clef == "size" : size = valeur
            if clef == "font": font = valeur
            if clef == "prefix" : prefix = valeur
            if clef == "suffix" : suffix = valeur
            if clef == "type_date" : type_date = valeur

#-------------------------------------------------------------------------------------------------------------  
    if prefix: PPTX_add_run(le_paragraph,f'{prefix} ',
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)
            
    prenom = text_personne(MH_personne,"prenom")
    prenoms = text_personne(MH_personne,"prenoms")
    nom = text_personne(MH_personne,"nom")
    surnom= text_personne(MH_personne,"surnom")
    conjugaison = "e" if text_personne(MH_personne,"sexe") == "F" else ""
    bdate = text_personne(MH_personne,"bdate")
    byear = text_personne(MH_personne,"byear")
    bplace = text_personne(MH_personne,"bplace")
    bcity = text_personne(MH_personne,"bcity")
    ddate = text_personne(MH_personne,"ddate")
    dyear = text_personne(MH_personne,"dyear")
    dplace = text_personne(MH_personne,"dplace")
    dcity = text_personne(MH_personne,"dcity")
    lacause = text_personne(MH_personne,"lacause")
    bdyear = text_personne(MH_personne,"bdyear")

    if prenom and isPrenom:   
        PPTX_add_run(le_paragraph,f'{prenom}', 
                            "bold",
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)
    if nom and isNom:         
        PPTX_add_run(le_paragraph,f' {nom}', 
                            "bold",
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)
        
    if prenoms and isPrenoms: 
        PPTX_add_run(le_paragraph,f' {prenoms}',
                            "italic",
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)
    if surnom and isSurnom:   
        PPTX_add_run(le_paragraph,f' "{surnom}"',
                            "italic",
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)

    if isYear or isDate : 

        if isYear:
            PPTX_add_run(le_paragraph,f' ({bdyear})',            #{byear}{"-"+dyear+"†" if dyear else ""})',
                                "italic",
                                "underline" if isUnderline else "",
                                size=size,
                                font=font,                         
                                color=Bleu)
            
            PPTX_add_run(le_paragraph,f' {" "+ bcity if isbLocation and bcity else ""}{" †"+ dcity[0] if (isdLocation and dcity) else ""}',
                                "italic",
                                "underline" if isUnderline else "",
                                size=size,
                                font=font,                         
                                color=Bleu)      
 
        if isDate:
            if type_date == 0:
                PPTX_add_run(le_paragraph,' (',size=size,font=font,color=color)
                PPTX_add_run(le_paragraph,f'{bdate}',
                                "underline" if isUnderline else "",size=size,font=font,color=Bleu)
                            
                PPTX_add_run(le_paragraph,f' à {bcity}' if bcity else "",
                                "underline" if isUnderline else "",size=size,font=font,color=Bleu)
                
                if ddate: 
                    PPTX_add_run(le_paragraph,f', † {ddate}',size=size,font=font,color=color)
                    if dplace :
                        PPTX_add_run(le_paragraph,f' à {dcity}',
                                "underline" if isUnderline else "",size=size,font=font,color=Bleu)
                    if lacause and isCause:
                        PPTX_add_run(le_paragraph,f' ,{lacause})',
                                "underline" if isUnderline else "",size=size,font=font,color=Bleu)
                PPTX_add_run(le_paragraph,')',size=size,font=font,color=color)

            # entourage
            elif type_date == 3:
                PPTX_add_run(le_paragraph,f', né{conjugaison} le {bdate}',
                                "underline" if isUnderline else "",size=size,font=font,color=color)
                                
                PPTX_add_run(le_paragraph,f' à {bcity}' if bplace else "",
                                "underline" if isUnderline else "", size=size,font=font,color=color)
            
                if ddate : 
                        PPTX_add_run(le_paragraph,f', décédé{conjugaison} le {ddate}',size=size,font=font,color=color)
                        if dplace :
                            PPTX_add_run(le_paragraph,f' à {dcity}',
                                    "underline" if isUnderline else "",
                                    size=size,
                                    font=font,
                                    color=color)
            # ascendant draw
            elif type_date == 2:
                    PPTX_add_run(le_paragraph,f'{bdate}',
                                    "underline" if isUnderline else "",size=size,font=font,color=color)

                    PPTX_add_run(le_paragraph,f' à {bcity}' if bplace else "",
                                    "underline" if isUnderline else "",size=size,font=font,color=color)
                    
                    if ddate : 
                        #PPTX_add_run(le_paragraph,f', † {ddate}',size=size,font=font,color=color)
                        PPTX_add_run(le_paragraph,f'\n† {ddate}',size=size,font=font,color=color)
                        if dplace :
                            PPTX_add_run(le_paragraph,f' à {dcity}',
                                    "underline" if isUnderline else "",size=size,font=font,color=color)

                        if lacause and isCause:
                            PPTX_add_run(le_paragraph,f' ,{lacause})',
                                    "underline" if isUnderline else "",size=size,font=font,color=color)              
            # descendants draw
            elif type_date == 4:
                PPTX_add_run(le_paragraph,f'{bdate}',
                                "underline" if isUnderline else "",size=size,font=font,color=color)

                #PPTX_add_run(le_paragraph,f' à {blocation[0]}' if blocation else "",
                #                "underline" if isUnderline else "",size=size,font=font,color=color)
                
                if ddate : 
                    #PPTX_add_run(le_paragraph,f', † {ddate}',size=size,font=font,color=color)
                    PPTX_add_run(le_paragraph,f' † {ddate}',size=size,font=font,color=color)
                    #if dlocation :
                    #    PPTX_add_run(le_paragraph,f' à {dlocation[0]}',
                    #            "underline" if isUnderline else "",size=size,font=font,color=color)

                PPTX_add_run(le_paragraph,f', {bcity}' if bplace else "",
                                "underline" if isUnderline else "",size=size,font=font,color=color)
                
            # descendants draw
            if type_date == 5:
                PPTX_add_run(le_paragraph,f' ({byear}{"-"+dyear+"†" if dyear else ""})',
                                    size=size,font=font,color=color)
                
            elif type_date == 6: # descendants_entourage
                PPTX_add_run(le_paragraph,f', {bdate}',size=size,font=font,color=color)
                PPTX_add_run(le_paragraph,f' à {bcity}' if bcity else "",size=size,font=font,color=color)
                
                if ddate : 
                    PPTX_add_run(le_paragraph,f'\n† {ddate}',size=size,font=font,color=color)
                    PPTX_add_run(le_paragraph,f' à {dcity}' if dcity else "",size=size,font=font,color=color)
            
     
# suffixe
    if suffix: PPTX_add_run(le_paragraph,f'{suffix}',
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)
#-------------------------------------------------------------------------------------------------------------   
    return 
#=============================================================================================================
def PPTX_add_paragraph_MH_personne(la_box,MH_personne,*args,**kwargs):
#-------------------------------------------------------------------------------------------------------------   
    if MH_personne == "???": return
    le_paragraph = None
#------------------------------------------------------------------------------------------------------------- 
    isCenter = False
    le_level = 0

    color = False
    size = False
    font= False
    prefix = False
    suffix= False

    isUnderline = False

    isNom = False
    isPrenom = False
    isPrenoms = False
    isSurnom = False
    isDate = False
    isYear = False
    isbLocation = False
    isdLocation = False
    isCause = False
    isBirth = False
    isDeath = False

    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
#lié au paragraph
            if valeur == "center" : isCenter = True                
# lié au run
            if valeur == "underline" : isUnderline =  True
# lié au contenu affiché
            if valeur == "nom" : isNom =  True
            if valeur == "prenom" : isPrenom =  True
            if valeur == "prenoms" : isPrenoms =  True
            if valeur == "surnom" : isSurnom =  True
            if valeur == "date" : isDate =  True
            if valeur == "year" : isYear =  True
            if valeur == "blocation" : isbLocation =  True
            if valeur == "dlocation" : isdLocation =  True
            if valeur == "birth" : isBirth =  True
            if valeur == "death" : isDeath = True
            if valeur == "dlocation" : isdLocation =  True
            if valeur == "cause" : isCause =  True

# lié au run
    for clef, valeur in kwargs.items():   
        if clef: 
            clef =clef.lower()
            if clef == "color" : color = valeur
            if clef == "size" : size = valeur
            if clef == "font": font = valeur
            if clef == "prefix" : prefix = valeur
            if clef == "suffix" : suffix = valeur
#lié au paragraph
            if clef == "level" : le_level = int(valeur)

#-------------------------------------------------------------------------------------------------------------   
    if prefix: 
        if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level )
        PPTX_add_run(le_paragraph,f'{prefix} ',
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)

    prenom,prenoms,nom,surnom,conjugaison = get_personne_ID(MH_personne)

    if prenom and isPrenom:   
        if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level )
        PPTX_add_run(le_paragraph,f'{prenom}', 
                            "bold",
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)
    if nom and isNom:   
        if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level )      
        PPTX_add_run(le_paragraph,f' {nom}', 
                            "bold",
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)
        
    if prenoms and isPrenoms:
        if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level ) 
        PPTX_add_run(le_paragraph,f' {prenoms}',
                            "italic",
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)
    if surnom and isSurnom:  
        if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level ) 
        PPTX_add_run(le_paragraph,f' "{surnom}"',
                            "italic",
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)

    if isYear or isDate or isBirth or isDeath: 
        bdate,byear,blocation = get_personne_info_naissance(MH_personne)
        ddate,dyear,dlocation,lacause = get_personne_info_deces(MH_personne)

    if isYear:
        if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level )
        PPTX_add_run(le_paragraph,f' ({byear}{"-"+dyear+"†" if dyear else ""})',
                            "italic",
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,                         
                            color=Bleu)
        
        PPTX_add_run(le_paragraph,f' {" "+ blocation[0] if isbLocation and blocation else ""}{" †"+ dlocation[0] if (isdLocation and dlocation) else ""}',
                            "italic",
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,                         
                            color=Bleu)      
    else:  
        if isDate :
            if bdate != "????":
                if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level )
                PPTX_add_run(le_paragraph,f', né{conjugaison} le {bdate}',
                                "underline" if isUnderline else "",
                                size=size,
                                font=font,
                                color=color)
                if blocation :        
                    PPTX_add_run(le_paragraph,f' à {blocation[0]}' if blocation else "",
                                    "underline" if isUnderline else "",
                                    size=size,
                                    font=font,
                                    color=color)

            if ddate : 
                if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level )
                PPTX_add_run(le_paragraph,f', décédé{conjugaison} le {ddate}',size=size,font=font,color=color)
                if dlocation :
                    PPTX_add_run(le_paragraph,f' à {dlocation[0]}',
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)

                if lacause and isCause:
                    if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level )
                    PPTX_add_run(le_paragraph,f' ({lacause})',
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)
                    
    if isBirth and ddate :
        if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level )
        PPTX_add_run(le_paragraph,f'* {bdate}',size=size,font=font,color=color)
        if blocation :
            PPTX_add_run(le_paragraph,f' ,{blocation[0]}',
                    "underline" if isUnderline else "",
                    size=size,
                    font=font,
                    color=color)
            
    if isDeath and ddate :
        if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level )
        PPTX_add_run(le_paragraph,f'† {ddate}',size=size,font=font,color=color)
        if dlocation :
            PPTX_add_run(le_paragraph,f' ,{dlocation[0]}',
                    "underline" if isUnderline else "",
                    size=size,
                    font=font,
                    color=color)

    if suffix: 
        if not le_paragraph: le_paragraph = PPTX_add_paragraph(la_box,"center" if isCenter else "", level = le_level )
        PPTX_add_run(le_paragraph,f'{suffix}',
                            "underline" if isUnderline else "",
                            size=size,
                            font=font,
                            color=color)
#-------------------------------------------------------------------------------------------------------------   
    return le_paragraph
#=============================================================================================================
def PPTX_add_couple_xy(sql_obj,la_page,x,y,box_width,box_height,images_style,MH_adult1,*args,**kwargs):
#------------------------------------------------------------------------------------------------------------- 
    type_date = 0
    la_couleur =""
    isAscendants = False
    le_style = "bold"
    isTransparent = False
    #
    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
            if valeur == "ascendants" : isAscendants = True   
            if valeur == "transparent" : isTransparent = True

    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "bcolor" : la_couleur = valeur
            if clef == "type_date" : type_date = valeur
            if clef == "style" : le_style = valeur

    #print(MH_adult1,MH_adult2)
    if MH_adult1 == MH_none : return

    if not la_couleur :
        if MH_adult1.sexe == "M": la_couleur =couleur_homme
        else: la_couleur = couleur_femme
            
    buffer_group = []

    image_size = box_height
    xb = x + image_size
    yb = y
    wb = box_width - image_size
    hb = box_height


# adult 1 
    le_texte_1 = text_personne(MH_adult1,"prenom","nom")

# photo adult 1
#WW1
    try:
        img = PPTX_add_photoID(sql_obj,la_page,x,y,image_size,image_size,MH_adult1)
        buffer_group.append(img)
    except Exception as error:
        print("PPTX_add_photoID : ",MH_adult1.name.format())

#-- Ajout du texte
    object_textbox = PPTX_add_box(la_page,xb ,yb,wb,hb,"transparent" if isTransparent else  "" ,bcolor = la_couleur,margin_left = 1,margin_right=1)
    buffer_group.append(object_textbox)
    #ligne 1 

# adult 1
    if le_texte_1 : 
        p = PPTX_add_paragraph(object_textbox,"center") #"left" if MH_adult2 != "???" else "center"
        PPTX_add_run(p,le_texte_1,le_style,box_width = wb)
        if type_date > 0:
            if type_date != 5 : p = PPTX_add_paragraph(object_textbox,"center")
            PPTX_add_run_MH_personne(p,MH_adult1,"date",type_date = type_date,size= 12 if isAscendants else 10)

#Grouping
    if buffer_group :
        le_groupe = la_page.shapes.add_group_shape(shapes=buffer_group) 

    return le_groupe
#=============================================================================================================
def PPTX_add_couple_xy_descendant(la_page,x,y,image_size,gap_immage,MH_adult1,MH_adult2,*args,**kwargs):
#-------------------------------------------------------------------------------------------------------------
# PPTX_add_couple_xy_descendant(la_page,x,y,image_size,MH_adult_1,MH_adult_2,"transparent" if isTransparent else "",bcolor = White,style="normal") 
    la_couleur =""
    isTransparent = False
    #
    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower() 
            if valeur == "transparent" : isTransparent = True

    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "bcolor" : la_couleur = valeur

#------------------------------------------------------------------------------------------------------------------------
    buffer_group = []
    if MH_adult1 != "???" : 

        # photo adult 1
        try:
            img = PPTX_add_photoID(la_page,x,y,image_size,image_size,MH_adult1)
            buffer_group.append(img)
        except Exception as error:
            print("PPTX_add_photoID : ",MH_adult1.name.format())

        # adult 2
        if MH_adult2 != "???":  

            try:
                img = PPTX_add_photoID(la_page,x+image_size+gap_immage,y,image_size,image_size,MH_adult2)
                buffer_group.append(img)

                tbox = PPTX_add_box(la_page,x+ image_size,y,gap_immage,image_size,*args,**kwargs)
                buffer_group.append(tbox)

            except Exception: print("PPTX_add_photoID : ",MH_adult2.name.format())

        #Grouping
        if buffer_group :le_groupe = la_page.shapes.add_group_shape(shapes=buffer_group) 

    return le_groupe
#=============================================================================================================
# PHOTOS
#=============================================================================================================
def PPTX_add_photoID(sql_obj,slide,x,y,img_height,img_width,MH_personne,*args):
#------------------------------------------------------------------------------------------------------------- 
    #if not get_personne_photoID_url(MH_personne) :return None
  
    isWatermark = False

    for valeur in args:  
        if valeur: 
            if isinstance(valeur,str): valeur = valeur.lower()
            if valeur == "watermark" : isWatermark = True

    img_path_temp = get_personne_photoID_file(sql_obj,MH_personne)

    if isWatermark: 
        w_prenom = MH_personne.prenom
        img_path = watermark(img_path_temp,w_prenom)
    else: img_path = img_path_temp

    sexe = MH_personne.sexe
#-- initiage image from local file
    img = Image.open(img_path)
#-- resize image
    r = img.height /img.width
    img_width = img_height / r
    if img_width > img_height:
        img_width = img_height

#-- add image in slide
    obj_image = slide.shapes.add_picture(img_path,Mm(x),Mm(y),Mm(img_width),Mm(img_height))
    if not isWatermark:
        if img_width < img_height : 
        #-- add gap in slide
            obj_gap = slide.shapes.add_textbox(Mm(x+img_width),Mm(y),Mm(img_height-img_width),Mm(img_height))
            if sexe == "M": couleur = couleur_homme
            else: couleur = couleur_femme
            obj_gap.fill.solid()
            obj_gap.fill.fore_color.rgb = hex_to_rgb(couleur)

        #-- Group
            le_groupe = slide.shapes.add_group_shape(shapes=[obj_image,obj_gap])  
            return le_groupe
        else:
            return obj_image
    else:
        return obj_image
#=============================================================================================================
def PPTX_add_photoIDs_side(sql_obj,la_page,temp_photoIDs,*args,**kwargs): 
#------------------------------------------------------------------------------------------------------------- 
    isDetails = False
    isImage = False
    isVertical = False
    y_margin_top = 0
    y_gap = 1
    box_width = 70
    x_margin_left = 0
    y_margin_bottom = 0
    y_margin_top = 0
    style = "bold"

    for valeur in args:
        if isinstance(valeur, str): valeur =valeur.lower()
        if valeur == "details" : isDetails = True
        if valeur == "image" : isImage = True
        
    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "x_margin_left" : x_margin_left = valeur
            if clef == "y_margin_top" : y_margin_top = valeur
            if clef == "box_width" : box_width = valeur
            if clef == "style" : style = valeur

#-------------------------------------------------------------------------------------------------------------
    ref_PhotoID = []
    les_photoIDs = []
    for MH_item in temp_photoIDs:
        if MH_item != MH_none:
            if get_personne_photoID_file(sql_obj,MH_item) and MH_item.indi_id not in ref_PhotoID:
                    les_photoIDs.append(MH_item)
                    ref_PhotoID.append(MH_item.indi_id)
#-------------------------------------------------------------------------------------------------------------  
    if les_photoIDs:
        nline = len(les_photoIDs)
        #la_box = PPTX_add_box(la_page,x_margin_left,y_margin_top,box_width,slide_height-y_margin_top,bcolor = GraySide)

        #calcul de la hauteur de la photo
        y = y_margin_top
        x = x_margin_left
        if isImage or isVertical:
            image_size = box_width
        else: image_size = min ((slide_height - y_margin_top - y_margin_bottom - (nline-1)*y_gap )/nline, 40)
        
        for MH_item in les_photoIDs:
            x = 0
            # photoID
            buffer_group = []
            try: img = PPTX_add_photoID(sql_obj,la_page,x,y,image_size,image_size,MH_item)
            except Exception as error: print("PPTX_add_photoIDs_side : ",text_personne_full(MH_item),error)

            # boite et nom
            if not isImage:

                if isVertical :              
                    xb = 0
                    yb = 70
                    wb = 70
                    hb = 25  
                else:   
                    xb = x + image_size
                    yb = y
                    wb = box_width - image_size
                    hb = image_size  

                # Ajout de la boite
                object_textbox = PPTX_add_box(la_page,xb,yb,wb,hb,"transparent",bcolor = White,margin_left = 1,margin_right=1)
                if not isImage : buffer_group.append(object_textbox)

                # Ajout des nomes prenoms
                le_prenom   = MH_item.prenom
                le_nom      = MH_item.nom

                if le_prenom : 
                    p = PPTX_add_paragraph(object_textbox,"center")
                    PPTX_add_run(p,le_prenom,style,box_width=wb)
                
                if le_nom : 
                    if image_size > 10:
                        p = PPTX_add_paragraph(object_textbox,"center")
                        PPTX_add_run(p,le_nom,style,box_width=wb)
                    else: PPTX_add_run(p,f' {le_nom}',style,box_width=wb)
                
                # Ajout des dates
                if isDetails:
                    PPTX_add_paragraph_MH_personne(object_textbox,MH_item,"prenoms","center",size = 10)
                    PPTX_add_paragraph_MH_personne(object_textbox,MH_item,"surnom","center",size = 10)
                    PPTX_add_paragraph_MH_personne(object_textbox,MH_item,"year","center",size = 10)

                # Grouping
                if buffer_group : le_groupe = la_page.shapes.add_group_shape(shapes=buffer_group) 
                # image suivante
                y = y + image_size + y_gap
            # cas de l'image seule
            else: break
#-------------------------------------------------------------------------------------------------------------                
    return
#=============================================================================================================
def PPTX_add_photoIDs_couple_side(la_page,temp_photoIDs,**kwargs): 
#-------------------------------------------------------------------------------------------------------------        
    y_margin_top = 0
    y_margin_bottom = 0
    y_gap = 1
    box_width = 70
    x_margin_left = 0
    style = "normal"
    image_size_max = 15
    
    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "x_margin_left" : x_margin_left = valeur
            if clef == "y_margin_top" : y_margin_top = valeur
            if clef == "y_margin_bottom" : y_margin_bottom = valeur
            if clef == "y_gap" : y_gap = valeur
            if clef == "box_width" : box_width = valeur
            if clef == "style" : style = valeur
            if clef == "image_size_max" : image_size_max = valeur
#-------------------------------------------------------------------------------------------------------------

    if not temp_photoIDs : return
    ref_PhotoID = []
    les_photoIDs = []
    for item in temp_photoIDs:
        ref0 = item[0].indi_id if item[0] != "???" else "???"
        ref1 = item[1].indi_id if item[1] != "???" else "???"
        if f'{ref0}{ref1}' and f'{ref1}{ref0}' not in ref_PhotoID and get_personne_photoID_url(item[0]):
            les_photoIDs.append(item)
            ref_PhotoID.append(f'{ref0}{ref1}')
            ref_PhotoID.append(f'{ref1}{ref0}')

#-------------------------------------------------------------------------------------------------------------  
    if les_photoIDs:
        PPTX_add_box(la_page,x_margin_left,y_margin_top,box_width,slide_height-y_margin_top,bcolor = GraySide)
#-------------------------------------------------------------------------------------------------------------               
        # calcul de la hauteur de la photo
        nline = len(les_photoIDs)
        y = y_margin_top
        image_size = min ((slide_height - y_margin_top - y_margin_bottom - (nline-1)*y_gap )/nline, image_size_max)
#-------------------------------------------------------------------------------------------------------------  
        for MH_item in les_photoIDs:
# boite et nom
            xb = x_margin_left + image_size
            yb = y
            hb = image_size  
            wb = box_width - 2 * image_size
# photoID
            buffer_group = []
            try:
                img = PPTX_add_photoID(la_page,x_margin_left,y,image_size,image_size,MH_item[0])
                if img : buffer_group.append(img)
            except Exception as error: print("PPTX_descendants_list_details : ",text_personne_full(MH_item[0]),error)

            if MH_item[1] != "???":
                try:
                    img = PPTX_add_photoID(la_page,x_margin_left + box_width-image_size,y,image_size,image_size,MH_item[1])
                    if img : buffer_group.append(img)
                except Exception as error:
                    print("PPTX_descendants_list_details : ",text_personne_full(MH_item[1]))  
            else:
                object_textbox = PPTX_add_box(la_page,box_width-image_size,yb,image_size,image_size,"transparent",bcolor = White,margin_left = 1,margin_right=1)

            # Ajout de la boite
            object_textbox = PPTX_add_box(la_page,xb,yb,wb,hb,"transparent",bcolor = White,margin_left = 1,margin_right=1)
            buffer_group.append(object_textbox)

            # Ajout des noms prenoms
            t1 = text_personne((MH_item[0]),"prenom","nom")
            t2 = text_personne((MH_item[1]),"prenom","nom")
            PPTX_add_run(PPTX_add_paragraph(object_textbox,"center"),t1,style,size=8 if " de " in t1 else 11)

            if MH_item[1] != "???": PPTX_add_run(PPTX_add_paragraph(object_textbox,"center"),t2,size=8 if " de " in t2 else 11)  

            # Grouping
            if buffer_group :le_groupe = la_page.shapes.add_group_shape(shapes=buffer_group) 

            # image suivante
            y = y + image_size + y_gap
    return
#=============================================================================================================
def PPTX_personne_MHphotos(sql_obj,le_document,MH_personne):
#------------------------------------------------------------------------------------------------------------

    y_top = 15 
    y_gap = 2
    n_row_max = 2
    x_gap = 2 
    x_bandeau = x_gap
    photo_height = (slide_height - y_top)/n_row_max - y_gap
#-------------------------------------------------------------------------------------------------------------
    buffer_photos = get_personne_MHphotos(sql_obj,MH_personne)
    #buffer_photos.append([photo_title+"|"+photo_date,img_path,img.width,img.height,le_watermark]) 
#-------------------------------------------------------------------------------------------------------------
    if buffer_photos:

        isNewPage = True
        for buffer_photo in buffer_photos:

            y_diff = 0

            # get MH_photo
            img_path = buffer_photo[1]
            img_width = buffer_photo[2]
            img_height = buffer_photo[3]
            #img = Image.open(img_path)
            # resize image
            percent_crop_horizontal = 0.0
            percent_crop_vertical = 0.0

            r = img_width / img_height
            if img_height > photo_height : 
                img_height = photo_height
                img_width = photo_height * r 

            # Nouvelle page si necessaire
            #-------------------------------------------------------------------------
            if isNewPage:
                x = x_gap
                y = y_top
                
                la_page, boxes = PPTX_add_page(le_document,slide_layout_album)
                p0 = PPTX_add_paragraph(boxes[0])
                PPTX_add_run(p0,text_personne_full(MH_personne))

                isNewPage = False
                isNext_row = False
                n_row = 1
                    
            else: 
            # sinon on calcule la position de la photo à ajouter
            #-------------------------------------------------------------------------

                x_diff = int( slide_width - x  - img_width - x_bandeau)

                if x_diff  < 0 : 

                    if slide_width - x - x_gap - x_bandeau > 30 :
                        percent_crop_horizontal = -1 * (x_diff / img_width)
                        isNext_row = True
                    else:
                        isNext_row = False
                        PPTX_expand_picture(pic_previous,x_bandeau)
                        n_row = n_row + 1
                        
                        if n_row > 2 :
                            # New page
                            x = x_gap
                            y = y_top
                            
                            la_page, boxes = PPTX_add_page(le_document,slide_layout_album)
                            p0 = PPTX_add_paragraph(boxes[0])

                            isNewPage = False
                            isNext_row = False
                            n_row = 1

                        else :
                            x = x_gap
                            y = y + img_height + y_gap 

            # Ajout de l'image  
            # #------------------------------------------------------------------------- 
            xp = x

            if img_width > 0 and img_height > 0 :
                try:
                    pic = la_page.shapes.add_picture(img_path,Mm( xp ),Mm(y - y_diff),Mm(img_width),Mm(img_height))
                    pic.crop_left = percent_crop_horizontal
                    pic.crop_right = -1 * percent_crop_horizontal
                    pic.crop_top = percent_crop_vertical
                    pic.crop_bottom = -1 * percent_crop_vertical
                    pic_previous = pic
                except:
                    print("PPTX_personne_MHphotos:",text_personne_full(MH_personne))

            else:
                print("img_width=",img_width, "img_height=",img_height)
                exit()

            # page Recto
            if isNext_row :
                    n_row = n_row + 1
                    
                    if n_row > 2 :
                        isNewPage = True
                    else :
                        x = x_gap
                        y = y + photo_height + y_gap  
                        isNext_row = False
            else:

                x = x + img_width + x_gap 
    return
#=============================================================================================================
def PPTX_expand_picture(pic,x_bandeau) :
        if pic:
            pvh = pic.height
            r = pic.width / pvh

            pic.width = (slide_width - x_bandeau ) * EMU - pic.left 
            pic.height = int(pic.width / r)

            pic.crop_bottom = -(1-pvh/pic.height)
            pic.crop_top = (1 - pvh/pic.height)
        return
#=============================================================================================================
# NOTES
#=============================================================================================================
def PPTX_add_note_texte(HTML_notes,*args):
#------------------------------------------------------------------------------------------------------------- 
    if not HTML_notes : return 
    onlyfirst = False
    for value in args:
        if value == "first" : onlyfirst = True
#------------------------------------------------------------------------------------------------------------- 
    class MyHTMLParser(HTMLParser):
        def handle_starttag(self, tag, attrs):
            la_liste = ["starttag"]
            la_liste.append(tag)
            for attr in attrs: la_liste.append(attr)
            les_data.append(la_liste)
        def handle_endtag(self, tag):
            les_data.append(["endtag",tag])
        def handle_data(self, data):
            les_data.append(["data",data])
    parser = MyHTMLParser()
    les_data = []
    isFirst = True
    is_linkurl = False
    is_linkname = False
    le_texte = ""
#------------------------------------------------------------------------------------------------------------- 
    for HTLM_note in HTML_notes:
        if HTLM_note == "Current occupation:1": le_texte = "Situation actuelle"
        else :
            parser.feed(HTLM_note)

            for la_data in les_data:
        # starttag 
                if la_data[0] == "starttag":
                    if la_data[1] == "linkurl"  : is_linkurl = True
                    if la_data[1] == "linkname" : is_linkname = True
                    if la_data[1] == "strong": le_texte = le_texte + "<strong>"
        #data
                if la_data[0] == "data":
                    if la_data[1] != "Web content link:" and la_data[1] != " " and not is_linkurl and not is_linkname:
                        if isFirst: 
                            isFirst = False
                            le_texte =  f'{la_data[1]}'
                        else: 
                            le_texte = f'{le_texte}, {la_data[1]}'
        #--- endtags
                if la_data[0] == "endtag":
                    if la_data[1] == "linkurl" : is_linkurl = False
                    if la_data[1] == "linkname" : is_linkname = False
                    if la_data[1] == "strong": le_texte = le_texte + "</strong>"

        if le_texte and onlyfirst : break
#-------------------------------------------------------------------------------------------------------------  
    return le_texte
#=============================================================================================================
def PPTX_add_run_html(p,HTML_text,*args,**kwargs):
#-------------------------------------------------------------------------------------------------------------
    isForcedBold = False
    for cle in args:
        if cle == "bold" : 
            isForcedBold = True
#-------------------------------------------------------------------------------------------------------------
    class MyHTMLParser(HTMLParser):
        def handle_starttag(self, tag, attrs):
            la_liste = ["starttag"]
            la_liste.append(tag)
            for attr in attrs: la_liste.append(attr)
            les_data.append(la_liste)
        def handle_endtag(self, tag):
            les_data.append(["endtag",tag])
        def handle_data(self, data):
            les_data.append(["data",data])
#-------------------------------------------------------------------------------------------------------------           
    parser = MyHTMLParser()
    les_data = []
    isBold = False
#------------------------------------------------------------------------------------------------------------- 
    if HTML_text :  
        #print(HTML_text)

        parser.feed(HTML_text)

        for la_data in les_data:
            #print(la_data)
            # starttag 
            if la_data[0] == "starttag":
                if la_data[1] == "strong"  : isBold = True
            # data
            if la_data[0] == "data":
                #print(la_data[1])
                PPTX_add_run(p,f'{la_data[1]}',"bold" if isBold or isForcedBold else "",*args,**kwargs)

            # endtags
            if la_data[0] == "endtag":
                if la_data[1] == "strong" : isBold = False
#-------------------------------------------------------------------------------------------------------------  
    return 
#=============================================================================================================
# ASCENDANT, DESCENDANT, ENTOURAGE
#=============================================================================================================
def PPTX_biographies_light_table(sql_obj,le_document,MH_personnes,box_sommaire):

    clefs_MH_personnes = []
    clean_list = []
    for item in MH_personnes:

        MH_personne = item[0]
        le_type = item[1] #ascendant,descendant
        MH_cible = item[2]
        n_level = item[3]
        la_Celebrity = item[4]

        les_bios = get_personne_bios(sql_obj,MH_personne)
        if not la_Celebrity and not les_bios :

            #calcul des lignées pour la cible
            if le_type == "ascendant":
                les_lignées = get_personne_lignées(sql_obj,MH_cible)
                la_lignée,la_lignée_idx = get_personne_lignée(MH_personne,les_lignées)
                temp_lignée = []
                for idx,MH_indi in enumerate(la_lignée):
                    temp_lignée.append(f'{f'[{idx}] ' if idx > 0 else ""}{text_personne(MH_indi,"prenom","nom","bdyear")}')
                la_lignée_texte = " < ".join(temp_lignée)
            else: 
                la_lignée_texte = ""
                la_lignée_idx = 0
            
            clef = f'{MH_personne.indi_id}{le_type}{MH_cible.indi_id}{n_level}'
            if clef not in clefs_MH_personnes and not isinstance(item[0],str): 
                clefs_MH_personnes.append(clef)
                clean_list.append([MH_personne,f'{n_level}|{MH_personne.nom.upper()}',la_lignée_texte,la_lignée_idx,
                                   le_type,f'{MH_cible.nom} {MH_cible.prenom} {MH_cible.prenoms}'])
    
    if clean_list:
        clefs_tri = list_unique_colomn(clean_list,4)
        clefs_cible = list_unique_colomn(clean_list,5)  

        ml = 10
        c = [20,11,45,36,13,13,0]
        c[-1] = slide_width-ml - sum(c[:-1])

        col_header_format = [
                    [c[0],"center","normal",14,Black,GraySide,Black,GraySide,1,1,1,1], 
                    [slide_width-ml-sum(c[:-1]),"center","normal",14,Black,GraySide,Black,GraySide,1,1,1,1], 
                    ]

        col_format = [
                    [c[0],"center","normal",14,White,couleur_chemin,White,couleur_chemin,1,1,1,1], 
                    [c[1],"center","normal",14,White,couleur_chemin,White,couleur_chemin,1,1,1,1], 
                    [c[2],"left","bold",14,White,couleur_chemin,White,couleur_chemin,2,1,1,1],
                    [c[3],"left","normal",14,White,couleur_chemin,White,couleur_chemin,1,1,1,1],
                    [c[4],"center","normal",14,White,Darkblue,White,Darkblue,0.5,0.5,1,1],
                    [c[5],"center","normal",14,White,Darkblue,White,Darkblue,0.5,0.5,1,1],
                    [c[6],"left","normal",14,White,Darkblue,White,Darkblue,2,1,1,1]
                    ]
        
        col_merged_format =[
                    [c[0],"center","normal",14,Black,White,Black,White,1,1,1,1], 
                    [c[1],"center","normal",14,Black,White,Black,GrayRow,1,1,1,1], 
                    [slide_width-ml-sum(c[:-1]),"left","normal",14,Black,White,Black,GrayRow,2,1,1,1]
                    ]
        
        col_lignée_format =[
                    [c[0],"center","normal",14,Black,White,Black,White,1,1,1,1], 
                    [c[1],"center","bold",14,Black,White,Black,GrayRow,1,1,1,1], 
                    [slide_width-ml-sum(c[:-1]),"left","normal",14,couleur_chemin,White,couleur_chemin,GrayRow,2,1,1,1]
                    ]
        
        col_href_format =[
                [c[0],"center","normal",14,Black,White,Black,White,1,1,1,1], 
                [c[1],"center","normal",14,Darkblue,White,Darkblue,White,1,1,1,1], 
                [slide_width-ml-sum(c[:-1]),"left","normal",14,Darkblue,White,Darkblue,GrayRow,2,1,1,1]
                ]
            
        max_n_col = len(col_format)-1
    #--------------------------------------------------------------------------------------              
        for clef_tri in clefs_tri:
            for clef_cible in clefs_cible:
    #---------------------------------------------------------------------------------------         
                filter_list = filter(lambda c: (c[4]==clef_tri and c[5]== clef_cible),clean_list)
                clean_list = sorted(filter_list, key=lambda col: (col[1]))
                if  clean_list : 
                
                    display_table = []
                    isNew = True

                    if clef_tri == "cible": le_titre = f"Biographie"
                    else: le_titre = f"Biographies des {clef_tri}s"
                    
                    for clean_item in clean_list:

                        MH_item  = clean_item[0] 
                        la_lignée_texte = clean_item[2] 
                        la_lignée_idx = clean_item[3] 
                        la_cible = clean_item[5]

                        le_header = [isNew,True,max_n_col,
                                    [col_header_format[0] + [" "],
                                     col_header_format[1] + [f'{le_titre} de {la_cible}']],0]

                        les_educs = get_personne_events(sql_obj,MH_item,even ="EDUC")
                        les_hrefs = get_personne_hrefs(sql_obj,MH_item)
                        les_occus = les_occus = get_personne_events(sql_obj,MH_item,even="OCCU")
                        les_maisons = get_personne_events(sql_obj,MH_item,even="CENS")
                        les_events = get_personne_events(sql_obj,MH_item,even="EVEN")

                        if les_occus or les_educs or les_events or les_maisons or les_hrefs:

                            prenom = text_personne(MH_item,"prenom")
                            prenoms = text_personne(MH_item,"prenoms")
                            nom = text_personne(MH_item,"nom")
                            surnom = text_personne(MH_item,"surnom")
                            byear = text_personne(MH_item,"byear")
                            dyear = text_personne(MH_item,"dyear")

                            display_table.append( 
                                                [isNew,False,max_n_col,
                                                    [
                                                        col_format[0] + [MH_item], 
                                                        col_format[1] + [" "],
                                                        col_format[2] + [nom],
                                                        col_format[3] + [f'{prenom if prenom else ""}'],
                                                        col_format[4] + [f'{byear if byear else "????"}'],
                                                        col_format[5] + [f'{dyear if dyear else ""}'],
                                                        col_format[6] + [f'{prenoms if prenoms else ""} {surnom if surnom else ""}']
                                                    ] ,0])     
                            isNew = False
                            isEven = True

                            if la_lignée_texte : 
                                if la_lignée_idx > 0:

                                    texte_chapitre = "👨‍👨‍👧‍👦"
                                    display_table.append([False,isEven, max_n_col,      
                                                            [col_lignée_format[0] + [None],
                                                             col_lignée_format[1] + [texte_chapitre],
                                                             col_lignée_format[2] + [[f'Ancètre direct de {clean_item[3]}{"ème" if la_lignée_idx> 1 else "er"} génération']]]
                                                            ])
                                    display_table.append([False,isEven, max_n_col,      
                                                            [col_lignée_format[0] + [None],
                                                             col_lignée_format[1] + [texte_chapitre],
                                                             col_lignée_format[2] + [la_lignée_texte]],0
                                                            ])
                                                          
                                    isEven = False if isEven else True
                                    isNew = False

                            if les_maisons:
                                for idx,item in enumerate(les_maisons):
                                    le_texte = ""
                                    if item.date  : le_texte = f'<STRONG>{item.date}</STRONG>: ' 
                                    if item.place : le_texte = f'{le_texte} {item.place}'   

                                    display_table.append([isNew,isEven, max_n_col, 
                                                    [col_lignée_format[0] + [None],              
                                                     col_merged_format[1] + ["🏠" if idx == 0  else " "],
                                                     col_merged_format[2] + [le_texte]],0
                                                    ]) 
                                    isNew = False

                                isEven = False if isEven else True
            
                            if les_educs:
                                for idx,item in enumerate(les_educs):
                                    le_texte = ""
                                    if item.date : le_texte = f'<STRONG>{item.date}</STRONG>: '
                                    if item.description : le_texte = f'{le_texte} {item.description}' 
                                    if item.place : le_texte = f'{le_texte} à {item.place}'
                                    if item.note : le_texte = f'{le_texte}: {item.note}'
                                    
                                    display_table.append([isNew,isEven,max_n_col, 
                                                    [col_lignée_format[0] + [None],           
                                                     col_merged_format[1] + ["🗞️" if idx == 0 else " "],
                                                     col_merged_format[2] + [le_texte]
                                                    ],0] )  
                                    isNew = False                      
                                    
                                isEven = False if isEven else True

                            if les_occus:
                                for idx,item in enumerate(les_occus):
                                    le_texte = ""
                                    if item.date : le_texte = f'<STRONG>{item.date}</STRONG>: '
                                    if item.description : le_texte = f'{le_texte} {item.description}' 
                                    if item.place : le_texte = f'{le_texte} à {item.place}'
                                    if item.note : le_texte = f'{le_texte}: {item.note}'

                                    display_table.append(   
                                                [isNew,isEven,max_n_col, 
                                                    [col_lignée_format[0] + [None],          
                                                     col_merged_format[1] + ["🛠️" if idx == 0  else " "],
                                                     col_merged_format[2] + [le_texte]
                                                    ],0])
                                    isNew = False   

                                isEven = False if isEven else True
                                                                                
                            if les_events:
                                for item in les_events:
                                    if item.type != "Celebrity":
                                        le_texte = ""
                                        if item.type : le_texte = f'{le_texte}<STRONG>{item.type}</STRONG>' 
                                        if item.description : le_texte = f'{le_texte}: {item.description}' 
                                        if item.date : le_texte = f'{le_texte}, {item.date}' 
                                        if item.place : le_texte = f'{le_texte} à {item.place}' 
                                        
                                        display_table.append([isNew,isEven, max_n_col,          
                                                        [ col_lignée_format[0] + [None], 
                                                          col_merged_format[1] + ["⭐"],
                                                          col_merged_format[2] + [le_texte]
                                                        ],0] ) 
                                        isNew = False

                                isEven = False if isEven else True
                                
                            if les_hrefs:
                                    
                                for idx,le_texte in enumerate(les_hrefs):
                                    display_table.append([isNew,isEven,max_n_col,      
                                                    [
                                                        col_lignée_format[0] + [None], 
                                                        col_href_format[1] + ["∞"],
                                                        col_href_format[2] + [le_texte]
                                                    ],0] ) 
                                    isNew = False
                            
                                isEven = False if isEven else True
                                           
                    if display_table:
        #------------------------------------------------------------------------------------------------------------------
                        PPTX_table_display_photoID_in_cell(sql_obj,le_document,col_format,le_header,display_table,box_sommaire,le_titre,
                                                           x_table=ml,image_width=c[0])

    return
#=============================================================================================================
def PPTX_biographies_table(sql_obj,le_document,MH_personnes,le_mode,box_sommaire):
#MH_personnes = list(0:MH_personne,1:type,2:MH_cible,3:level,4:Celebrity]
#ZOB
    clefs_MH_personnes = []
    clean_list = []
    
    for item in MH_personnes:
        MH_personne = item[0]
        le_type = item[1] #ascendant,descendant
        MH_cible = item[2]
        n_level = item[3]
        la_Celebrity = item[4]

        les_bios  = get_personne_bios(sql_obj,MH_personne)
        
        if ( (le_mode == "main") or
             (la_Celebrity and le_mode == "celebrity") or 
             (le_mode == "bio" and les_bios and not la_Celebrity)
           ):

            if le_type == "ascendant":
                les_lignées = get_personne_lignées(sql_obj,MH_cible)
                la_lignée,la_lignée_idx = get_personne_lignée(MH_personne,les_lignées)
                temp_lignée = []
                for idx,MH_indi in enumerate(la_lignée):
                    temp_lignée.append(f'{f'[{idx}] ' if idx > 0 else ""}{text_personne(MH_indi,"prenom","nom","bdyear")}')
                la_lignée_texte = " < ".join(temp_lignée)
            else: 
                la_lignée_texte = ""
                la_lignée_idx = 0

            clef = f'{MH_personne.indi_id}{le_type}{MH_cible.indi_id}{n_level}'
            if clef not in clefs_MH_personnes and not isinstance(item[0],str): 
                clefs_MH_personnes.append(clef)
                clean_list.append([MH_personne,f'{n_level}|{MH_personne.nom.upper()}',la_lignée_texte,la_lignée_idx])
            
    if clean_list : 

        clean_list = sorted(clean_list, key=lambda col: (col[1]))  

        # parametrage de la table
        ml = 30
        c = [11,65,50,14,14,0]
        c[-1] = slide_width-ml - sum(c[:-1])
        sl = 16
        col_format = [
                    [c[0],"center","normal",sl,White,couleur_chemin,White,couleur_chemin,1,1,1,1], 
                    [c[1],"left","bold",sl,White,couleur_chemin,White,couleur_chemin,2,1,1,1],
                    [c[2],"left","normal",sl,White,couleur_chemin,White,couleur_chemin,1,1,1,1],
                    [c[3],"center","normal",sl,White,Darkblue,White,Darkblue,0.5,0.5,1,1],
                    [c[4],"center","normal",sl,White,Darkblue,White,Darkblue,0.5,0.5,1,1],
                    [c[5],"left","normal",sl,White,Darkblue,White,Darkblue,2,1,1,1]
                    ]
                    
        col_merged_format =[
                    [c[0],"center","normal",14,Black,White,Black,White,1,1,1,1], 
                    [slide_width-21,"left","normal",14,Black,White,Black,GrayRow,2,1,1,1]
                    ]
        
        col_lignée_format =[
                    [c[0],"center","normal",14,Black,White,Black,White,1,1,1,1], 
                    [slide_width-21,"left","normal",14,couleur_chemin,White,couleur_chemin,GrayRow,2,1,1,1]
                    ]
        
        col_lignée_bold_format =[
                    [c[0],"center","bold",14,Black,White,Black,White,1,1,1,1], 
                    [slide_width-21,"left","bold",14,couleur_chemin,White,couleur_chemin,GrayRow,2,1,1,1]
                    ]
            
        col_href_format =[
                    [c[0],"center","normal",14,Darkblue,White,Darkblue,White,1,1,1,1], 
                    [slide_width-21,"left","normal",14,Darkblue,White,Darkblue,GrayRow,2,1,1,1]
                    ]
        
        col_href_bold_format =[
                    [c[0],"center","bold",14,Darkblue,White,Darkblue,White,1,1,1,1], 
                    [slide_width-21,"left","bold",14,Darkblue,White,Darkblue,GrayRow,2,1,1,1]
                    ]
        max_n_col = len(col_format)-1
#-------------------------------------------------------------------------------------- 
        
        for clean_item in clean_list:

            display_table = []
            isNew = True
            MH_item  = clean_item[0] 
            la_lignée_texte = clean_item[2] 
            la_lignée_idx = clean_item[3] 

            les_educs = get_personne_events(sql_obj,MH_item,even ="EDUC")
            if le_mode != "celebrity" : les_maisons = get_personne_events(sql_obj,MH_item,even="CENS")
            else: les_maisons = []

            if le_mode != "celebrity" : les_occus = get_personne_events(sql_obj,MH_item,even="OCCU")
            else: les_occus = []

            les_events = get_personne_events(sql_obj,MH_item,even = "EVEN")            
            les_bios  = get_personne_bios(sql_obj,MH_item)
            les_hrefs = get_personne_hrefs(sql_obj,MH_item)

            if les_occus or les_educs or les_events or les_bios or les_hrefs or les_maisons:

                prenom = text_personne(MH_item,"prenom")
                prenoms = text_personne(MH_item,"prenoms")
                nom = text_personne(MH_item,"nom")
                surnom = text_personne(MH_item,"surnom")
                byear = text_personne(MH_item,"byear")
                dyear = text_personne(MH_item,"dyear")
                
                le_header = [True, False ,max_n_col,[
                                col_format[0] + [" "],
                                col_format[1] + [nom],
                                col_format[2] + [f'{prenom if prenom else ""}'],
                                col_format[3] + [f'{byear if byear else "????"}'],
                                col_format[4] + [f'{dyear if dyear else ""}'],
                                col_format[5] + [f'{prenoms if prenoms else ""} {surnom if surnom else ""}']
                            ],0,MH_item] 
              
                isEven = True

                if la_lignée_texte : 
                    if la_lignée_idx > 0:

                        display_table.append([isNew,isEven,max_n_col,          
                                            [col_lignée_bold_format[0] + ["👨‍👨‍👧‍👦"],
                                            col_lignée_bold_format[1] + [f'Ancètre direct de {la_lignée_idx}{"ème" if la_lignée_idx> 1 else "er"} génération']
                                            ], 0,MH_item] ) 
                        isEven = False if isEven else True
                        isNew = False

                        display_table.append([isNew,isEven,max_n_col,          
                                            [col_lignée_format[0] + ["👨‍👨‍👧‍👦"],
                                            col_lignée_format[1] + [la_lignée_texte]
                                            ], 0,MH_item] ) 
                        
                        isEven = False if isEven else True
                        isNew = False

                if les_maisons:
                    for idx,item in enumerate(les_maisons):
                        le_texte = ""
                        if item.date  : le_texte = f'<STRONG>{item.date}</STRONG>: ' 
                        if item.place : le_texte = f'{le_texte} {item.place}'   

                        display_table.append(   
                                    [isNew,isEven, max_n_col,          
                                        [col_merged_format[0] + ["🏠" if idx == 0  else " "],
                                            col_merged_format[1] + [le_texte]
                                        ],0,None]) 
                        isNew = False

                    isEven = False if isEven else True
   
                if les_educs:
                    for idx,item in enumerate(les_educs):
                        le_texte = ""
                        if item.date : le_texte = f'<STRONG>{item.date}</STRONG>: '
                        if item.description : le_texte = f'{le_texte} {item.description}' 
                        if item.place : le_texte = f'{le_texte} à {item.place}'
                        if item.note : le_texte = f'{le_texte}: {item.note}'
                        
                        display_table.append([isNew,isEven,max_n_col,           
                                        [col_merged_format[0] + ["🗞️" if idx == 0 else " "],
                                         col_merged_format[1] + [le_texte]
                                        ], 0, MH_item] )  
                        isNew = False                      
                        
                    isEven = False if isEven else True

                if les_occus:
                    for idx,item in enumerate(les_occus):
                        le_texte = ""
                        if item.date : le_texte = f'<STRONG>{item.date}</STRONG>: '
                        if item.description : le_texte = f'{le_texte} {item.description}' 
                        if item.place : le_texte = f'{le_texte} à {item.place}'
                        if item.note : le_texte = f'{le_texte}: {item.note}'

                        display_table.append(   
                                    [isNew,isEven,max_n_col,          
                                        [col_merged_format[0] + ["🛠️" if idx == 0  else " "],
                                            col_merged_format[1] + [le_texte]
                                        ],0,None])
                        isNew = False   

                    isEven = False if isEven else True
                                                                      
                if les_events:
                    for item in les_events:
                        if item.type != "Celebrity":
                            le_texte = ""
                            if item.type : le_texte = f'{le_texte}<STRONG>{item.type}</STRONG>' 
                            if item.description : le_texte = f'{le_texte}: {item.description}' 
                            if item.date : le_texte = f'{le_texte}, {item.date}' 
                            if item.place : le_texte = f'{le_texte} à {item.place}' 
                            
                            display_table.append([isNew,isEven, max_n_col,          
                                            [ col_merged_format[0] + ["⭐"],
                                            col_merged_format[1] + [le_texte]
                                            ], 0 , MH_item] ) 
                            isNew = False

                    isEven = False if isEven else True
                    
                if les_bios:
                    for idx,le_texte in enumerate(les_bios):    
                        display_table.append([isNew,isEven,max_n_col,      
                                        [
                                            col_merged_format[0] + ["✒️"  if idx == 0  else " "],
                                            col_merged_format[1] + [le_texte]
                                        ], 0 , MH_item] ) 
                        isNew = False
                            
                    isEven = False if isEven else True
                        

                if les_hrefs:
   
                    display_table.append([isNew,isEven,max_n_col,      
                                    [
                                        col_href_bold_format[0] + ["∞"],
                                        col_href_bold_format[1] + ["Pour plus d'informations, suivez les liens"]
                                    ], 0, MH_item] ) 
                        
                    for idx,le_texte in enumerate(les_hrefs):
                        display_table.append([isNew,isEven,max_n_col,      
                                        [
                                            col_href_format[0] + [""],
                                            col_href_format[1] + [le_texte]
                                        ], 0 , MH_item] ) 
                        isNew = False
                
                    isEven = False if isEven else True
                                  
                                    
            if display_table:
#------------------------------------------------------------------------------------------------------------------
                le_titre = f'{text_personne(MH_item,"prenom","nom")}, {clean_item[1]}'
                PPTX_table_display(sql_obj,le_document,col_format,le_header,display_table,box_sommaire,x_table = ml )

                if le_mode == "celebrity" :  PPTX_personne_MHphotos(sql_obj,le_document,MH_item )

    return
#=============================================================================================================
def PPTX_add_ligne_sommaire(box_sommaire,le_document,la_page,le_texte,*args):

    if box_sommaire and le_texte : 

        isSection = False
        for valeur in args:  
            if valeur: 
                if isinstance(valeur,str): valeur =valeur.lower()
                if valeur == "section" : isSection = True

        if isinstance(le_texte, list): the_texte = " ".join(le_texte)
        elif isinstance(le_texte, str): the_texte = le_texte
        else : the_texte = "Erreur; pas de texte dans PPTX_add_ligne_sommaire"

        p = PPTX_add_paragraph(box_sommaire,"left",level = 1)
        slide_id = le_document.slides.index(la_page) + 1
        le_run = PPTX_add_run(p,f'Page {slide_id:03d}')
        if le_run : PPTX_add_hyperlink_to_slide(box_sommaire,le_run,la_page)

        if isSection :
            PPTX_add_run(p,f'  ')
            PPTX_add_run(p,f'{the_texte}',"bold")
        else:
            PPTX_add_run(p,f'  • ',size=12)
            PPTX_add_run(p,f'{the_texte}',size=12)

    return 
#=============================================================================================================
def PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_cible,direction,sujet,**kwargs):
#-------------------------------------------------------------------------------------------------------------
    isVerbose = False
#-------------------------------------------------------------------------------------------------------------
    n_level_max = 999
    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "n_level_max" : n_level_max = valeur
#-------------------------------------------------------------------------------------------------------------
    if direction == "ascendant":
        les_couples = get_personne_ascendants(sql_obj,MH_cible,[],0,n_level_max)
        #list(MH_couple({"level":n_level, "adult1":adult1 , "adult2":adult2, "bio":isAdopted })
        for le_couple in les_couples:
            if le_couple.level != 1 : 
                MH_personnes.append([le_couple.adult1,sujet,MH_cible,le_couple.level-1,get_personne_celebrity(sql_obj,le_couple.adult1)])
                if le_couple.adult2 != MH_none : 
                    MH_personnes.append([le_couple.adult2,sujet,MH_cible,le_couple.level-1,get_personne_celebrity(sql_obj,le_couple.adult2)])

    else: # descendant
        t = sujet
        if sujet == "couzmater" : t = "cousinades maternelle"
        elif sujet == "couzpater" : t = "cousinades paternelle"
        MH_items,le_level_max = get_personne_entourage(sql_obj,MH_cible,n_level_max,sujet)
        if MH_items : 
            for MH_item in MH_items:
                if MH_item[0] > 0:
                    if sujet != "descendant" or  MH_item[0] != 1 : 
                        MH_personnes.append([MH_item[1],t,MH_cible,MH_item[0]-1,get_personne_celebrity(sql_obj,MH_item[1])])
                        if MH_item[2] != MH_none : MH_personnes.append([MH_item[2],t,MH_cible,MH_item[0]-1,get_personne_celebrity(sql_obj,MH_item[2])])

    if isVerbose :
        for i in MH_personnes:
            print("personne_feed:",text_personne_full(i[0]),i[1],text_personne_full(i[2]),i[3],i[4])

    return MH_personnes
#=============================================================================================================      
def PPTX_descendant_arbre_table(sql_obj,le_document,MH_personne,couple_ref,les_sujets,n_level_max,box_sommaire,loption,image_size,**kwargs): 
#-------------------------------------------------------------------------------------------------------------
#ZQW
    # le titre  
    le_titre = f'Famille de {text_personne(MH_personne,"prenom")}'
    if couple_ref[0] != MH_none:
        le_titre =f'{le_titre} {"fille" if MH_personne.sexe =="F" else "fils"} de {text_personne(couple_ref[0],"prenom", "nom")} et {text_personne(couple_ref[1],"prenom", "nom")}'

    n_level_start = 3
    for clef, valeur in kwargs.items():
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "n_level_start" : n_level_start = valeur
#-------------------------------------------------------------------------------------------------------------  
    MH_entourages,le_level_max = get_personne_entourage(sql_obj,MH_personne,n_level_start,loption)
    # list(level, Adult1, Adult2, Année naissance)
    if MH_entourages:
        n_row_max = (slide_height - 11) / image_size

        if len(MH_entourages) > n_row_max :
            if n_level_start > 2 :
                PPTX_descendant_arbre_table(sql_obj,le_document,MH_personne,couple_ref,les_sujets,n_level_max,box_sommaire,loption,image_size,n_level_start = n_level_start-1)
                return
            else: 
                PPTX_descendant_arbre_table(sql_obj,le_document,MH_personne,couple_ref,les_sujets,n_level_max,box_sommaire,loption,image_size*0.9)
                return
        
        couple_ref = [MH_entourages[0][1],MH_entourages[0][2]]
#-------------------------------------------------------------------------------------------------------------
        ml = 10
        mb = 0.6
        mt = 0.6

        col_format_label = [
                    [slide_width-ml-image_size,"center","normal",14,Black,GraySide,Black,GraySide,1,1,1,1],
                    [20,"center","normal",14,Black,GraySide,Black,GraySide,1,1,1,1]
                    ]
        
        max_n_col = 0
        col_format_levels=[]

        if le_level_max > 1 : z_level = le_level_max
        else : z_level = 2

        for n_level in range(z_level):

            col_format_level = []  
            color_level = f'#{RGBColor(250-10*n_level,250-10*n_level,250-10*n_level)}'
        
            for k in range(0,2*(n_level)): 
                col_format_level.append([image_size,"center","normal",34,GraySide,White,GraySide,White,1,1,1,1]) 

            col_format_level.append([image_size,"center","normal",34,GraySide,color_level,GraySide,color_level,1,1,1,1]) 
            col_format_level.append([image_size,"center","normal",34,GraySide,color_level,GraySide,color_level,1,1,1,1]) 

            for k in range(2*(n_level),z_level):  
                col_format_level.append([image_size,"center","normal",16,Black,color_level,Black,color_level,1,1,1,1])

            col_format_level.append([slide_width-ml-(2*z_level)*image_size-10,"left","normal",14,Darkblue,color_level,Darkblue,color_level,2,1,mt,mb]) #noms et premons
            col_format_level.append([10,"center","bold",14,Darkblue,White,Darkblue,White,2,1,mt,mb])#  nenfant suite

            col_format_levels.append(col_format_level)

        max_n_col = len(col_format_level)-1

        le_header = [True,True,max_n_col,[col_format_label[0] + [le_titre]],0]
    #-------------------------------------------------------------------------------------------------------------

#-------------------------------------------------------------------------------------------------------------
        display_table = []
        isNew = True
        isEven = True
        suivants = []
        temp_suivants = []
        for MH_entourage in MH_entourages:

    # boucle sur chaque personne de l'entourage ---------------------------------------------------------------
            
            isEven = False if isEven else True

            n_level = MH_entourage[0]
            MH_adult1 = MH_entourage[1]
            MH_adult2 = MH_entourage[2]

            le_texte = text_couple(MH_adult1,MH_adult2)
            n_enfants = 0
            if n_level == le_level_max : 
                MH_enfants = get_couple_enfants(sql_obj,MH_adult1,MH_adult2)
                n_enfants = len(MH_enfants)

            la_display_ligne =[]
            for k in range(0,2*(n_level-1)): 
                la_display_ligne.append(col_format_levels[n_level-1][k] + [">" if k == 2*(n_level-1)-1 else None])

            la_display_ligne.append(col_format_levels[n_level-1][2*(n_level-1)] + [MH_adult1])
            la_display_ligne.append(col_format_levels[n_level-1][1+ 2*(n_level-1)] + [MH_adult2 if MH_adult2 != MH_none else None])

            for k in range(2*(z_level-n_level)): 
                la_display_ligne.append(col_format_levels[n_level-1][k+2+ 2*(n_level-1)] + [None])

            la_display_ligne.append(col_format_levels[n_level-1][-2] + [le_texte])
            la_display_ligne.append(col_format_levels[n_level-1][-1] + [f'> {n_enfants}' if n_enfants> 0 else None])
            display_table.append([isNew,isEven,max_n_col,la_display_ligne,0])            
            isNew = False

            if n_enfants > 0 : 
                if MH_adult1.indi_id not in temp_suivants:
                    temp_suivants.append(MH_adult1.indi_id)
                    suivants.append(MH_adult1)
                                                                                       
        if display_table:
    #------------------------------------------------------------------------------------------------------------------
            PPTX_table_display_photoID_in_cell(sql_obj,le_document,col_format_level,le_header,display_table,box_sommaire,le_titre,
                               x_table = ml ,layout = slide_layout_table_image )
    #------------------------------------------------------------------------------------------------------------------
        if suivants:
            for suivant in suivants:
                PPTX_descendant_arbre_table(sql_obj,le_document,suivant,couple_ref,les_sujets,n_level_max,box_sommaire,"descendant",20)

    return 
#=============================================================================================================
def PPTX_descendant_group_table(sql_obj,le_document,MH_personne,les_sujets,n_level_max,box_sommaire,loption,**kwargs): 
#-------------------------------------------------------------------------------------------------------------
#ZQS
    # le titre  
    le_titre = f'{les_sujets[0]}, {les_sujets[1]} et {les_sujets[2]} de {text_personne(MH_personne,"prenom")}'
    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "titre" : le_titre = valeur
#-------------------------------------------------------------------------------------------------------------
    ml = 20
    col_format_label =  [[slide_width-ml,"center","normal",16,Black,GraySide,Black,GraySide,1,1,1,1]]
    col_format_1 =      [[slide_width-ml,"center","normal",16,White,Darkblue,White,Darkblue,1,1,1,1]]
    col_format_ligne =  [[slide_width-ml,"center","normal",14,Darkblue,cell21,Darkblue,cell21,1,1,1,1]]
    
    max_n_col = len(col_format_1)-1
    le_header =     [True,False,max_n_col,[col_format_label[0] + [f'{le_titre}']],0,MH_personne]
#-------------------------------------------------------------------------------------------------------------
    MH_entourages,le_level_max = get_personne_entourage(sql_obj,MH_personne,n_level_max,loption)
    # list(level, adult1, adult2, Année naissance)
    if MH_entourages:
#-------------------------------------------------------------------------------------------------------------
        display_table = []
        isNew = True
        isEven = True
        # boucle sur chaque personne de l'entourage 
        for n_level in range(1,le_level_max+1):
            entourages_level = filter(lambda c: (c[0]==n_level),MH_entourages)
            entourages_level = sorted(entourages_level, key=lambda col: (col[3]) )
            # le titre
            if n_level < 4: t_titre = f'{les_sujets[n_level-1]}: {len(entourages_level)}' 
            else : t_titre = f'Niveau {n_level-1}: {len(entourages_level)}' 
            
            if n_level > 2 :isNew = True
            display_table.append( [isNew,False,max_n_col,[col_format_1[0] + [f'{t_titre}']],0,MH_personne])
            isNew =False

            for item in entourages_level:
                display_table.append( [isNew,isEven,max_n_col,[col_format_ligne[0] + [text_couple(item[1],item[2],"simple")]],0,MH_personne])
                isNew =False
                isEven = False if isEven else True
                                                          
        if display_table:
    #------------------------------------------------------------------------------------------------------------------
            PPTX_table_display(sql_obj,le_document,col_format_1,le_header,display_table,box_sommaire,x_table = ml)
    #------------------------------------------------------------------------------------------------------------------
    return 
#=============================================================================================================
def PPTX_ascendant_draw(sql_obj,le_document,MH_personne,n_col,le_sujet,box_sommaire,*args):  
    isSingle = False
    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
            if valeur == "single" : isSingle = True   
    #-------------------------------------------------------------------------------------------------------------
    MH_ascendants_max = get_personne_ascendants(sql_obj,MH_personne,[],0,n_col+1)
    #MH_ascendants = MH_couple({"level":n_level, "adult1":adult1 , "adult2":adult2, "bio":isAdopted }))
    #-------------------------------------------------------------------------------------------------------------
    if len(MH_ascendants_max) > 1: 
    #-------------------------------------------------------------------------------------------------------------
        #titres"

        le_titre = f'{" ".join(le_sujet)} de {text_personne(MH_personne,"prenom")}'
        la_page,boxes = PPTX_add_page(le_document,slide_layout_ascendants_4)
        PPTX_add_ligne_sommaire(box_sommaire,le_document,la_page,le_titre)
        PPTX_add_run(PPTX_add_paragraph(boxes[0]),le_titre,size = 18)

        # inititalisation des tailles
        n_row=1
        for x in range(n_col):
            n_row = n_row*2
        n_row = n_row + 1

        if n_col < 5 : gap_x = 10
        else:gap_x = 3
        gap_y = 1
        box_width = (slide_width - slide_margin_left - slide_margin_right - ((n_col-1) * gap_x) ) / n_col
        box_height = (slide_height - slide_margin_bottom - ((n_row-1) * gap_y) ) / n_row
        image_size = box_height
        #-------------------------------------------------------------------------------------------------------
        #Ajout des personnes
        for idx,couple in enumerate(MH_ascendants_max):
            MH_adult1 = couple.adult1
            if idx == 0:

                x0 = slide_margin_left                
                y0 = 0 

                x = x0
                y=y0

                i_row = 2
            else:       
                i_col = couple.level-1
                if couple.adult1.sexe == "F" : i_row = i_row + 1

                x = x0  + (i_col-1) * (box_width + gap_x )
                y = y0 + (i_row-1) * (box_height + gap_y)

                # ajout fleche gauche
                la_page.shapes.add_picture(icloud+'/MesProgrammes/MH_Photos/FLECHE_GAUCHE.png',
                                                        Mm(x-gap_x),Mm(y),Mm(gap_x),Mm(box_height))

            PPTX_add_couple_xy(sql_obj,la_page,x,y,box_width,box_height,image_size,MH_adult1,"ascendants",type_date = 2)
                                                        
#-------------------------------------------------------------------------------------------------------------
    return 
#=============================================================================================================
def PPTX_ascendant_table(sql_obj,le_document,MH_personne,les_lignées,n_level_max,n_rang,image_size):

    MH_ascendants = get_personne_ascendants(sql_obj,MH_personne,[],0,n_level_max) 
    #MH_couple({"level":n_level, "adult1":adult2 , "adult2":adult1, "bio":isAdopted })
    if MH_ascendants:

        ml = 10
        mb = 0.6
        mt = 0.6

        wtable = slide_width - ml
        wlevel = 10
        c = [image_size,10,0,16,16,40,40,30,10]
        c[2] = wtable - sum(c)

        max_n_col = 0
        col_format_levels=[]
        color_level =[]
        color_rang =[]

        if n_level_max > 1 : z_level = n_level_max - 1
        else : z_level = 2

        for n_level in range(z_level+1):
            color_level.append(f'#{RGBColor(250-10*n_level,250-10*n_level,250-10*n_level)}')
            color_rang.append(f'#{RGBColor(20,150-30*n_level,200-30*n_level)}')
           # "#EC5800"        

        col_header_format = [
                    [c[0],"center","normal",14,Black,GraySide,Black,GraySide,0,0,0,0], 
                    [wtable - c[0] ,"center","normal",14,Black,GraySide,Black,GraySide,0,0,0,0], 
                    ]
        
        col_lignée_format =[
            [c[0],"center","normal",14,Black,White,Black,White,1,1,1,1], 
            [wtable - c[0] - c[-1] ,"left","normal",14,Darkblue,White,Darkblue,White,0,0,1,1],
            [c[-1],"left","bold",20,White,White,White,White,2,1,0,0]#  suivant
            ]
        
        for n_level in range(z_level+1):
            la_couleur_level = color_level[n_level]
            la_couleur_rang = color_rang[n_level]

            col_format_level =[
            [c[0],"center","normal",14,Black,White,Black,White,1,1,1,1], # image
            [c[1],"center","bold",14,White,la_couleur_rang,White,la_couleur_rang,0,0,0,0] ] # rang
            #colonne niveau
            for k in range(0,n_level): 
                col_format_level.append([wlevel,"center","bold",20,la_couleur_rang,White,la_couleur_rang,White,0,0,0,0]) #niveau
            
            col_format_level.append([wtable - c[0] - c[1] - (n_level * wlevel) -sum(c[3:]) ,"left","normal",14,White,la_couleur_rang,White,la_couleur_rang,2,1,mt,mb]) #noms et premons
            col_format_level.append([c[3],"center","normal",14,Darkblue,la_couleur_level,Darkblue,la_couleur_level,0,0,0,0])#  byear
            col_format_level.append([c[4],"center","normal",14,Darkblue,la_couleur_level,Darkblue,la_couleur_level,0,0,0,0])#  dyear
            col_format_level.append([c[5],"center","normal",14,Darkblue,la_couleur_level,Darkblue,la_couleur_level,0,0,0,0])#  bcity
            col_format_level.append([c[6],"center","normal",14,Darkblue,la_couleur_level,Darkblue,la_couleur_level,0,0,0,0])#  bregion
            col_format_level.append([c[7],"center","normal",14,Darkblue,la_couleur_level,Darkblue,la_couleur_level,0,0,0,0])#  bcountry
            col_format_level.append([c[8],"left","bold",20,la_couleur_rang,White,la_couleur_rang,White,0,0,0,0])#  suivant
            
            col_format_levels.append(col_format_level)
            max_n_col = max(max_n_col,len(col_format_level))

        max_n_col = max_n_col - 1
        #print("max_n_col",max_n_col)
        
        #for  idx,t in enumerate(col_format_levels):
            #print (idx,"--"* 20)
            #for ii in t:
                #print (ii)
            
            #print("n_col",len(t),"-",max_n_col-1)

        #-------------------------------------------------------------------------------------- 
        #hauteur du header
        h_row = 0
        le_titre = f"ascendants de {text_personne_full(MH_personne)}"
        le_header = [True,True,max_n_col,[col_header_format[0] + [None],col_header_format[1] + [f'{le_titre}']],0]
        h_row = PPTX_table_row_height([col_header_format[1] + [f'{le_titre}']])

        la_lignée,la_lignée_idx = get_personne_lignée(MH_personne,les_lignées)
        if la_lignée:
            temp_lignée = []
            for MH_indi in la_lignée:
                temp_lignée.append(f'<STRONG>{text_personne(MH_indi,"prenom")} {text_personne(MH_indi,"nom")}</STRONG> {text_personne(MH_indi,"bdyear")}')
            la_lignée_texte = " < ".join(temp_lignée)

            h_row = h_row + PPTX_table_row_height([col_lignée_format[0] + [None],
                                            col_lignée_format[1] + [la_lignée_texte]])

        n_row_max = (slide_height - h_row) / image_size
        if len(MH_ascendants) > n_row_max :

            if image_size * 0.9  < 12 : 
                n_level_max = n_level_max - 1
            else : image_size = image_size * 0.9

            PPTX_ascendant_table(sql_obj,le_document,MH_personne,les_lignées,n_level_max,n_rang,image_size)
            return

        display_table = []
        isNew = True
        isEven = False

        if n_rang > 0:
            display_table.append([isNew,isEven, max_n_col,      
                                    [col_lignée_format[0] + [None],
                                     col_lignée_format[1] + [la_lignée_texte],
                                     col_lignée_format[2] + [" "]]
                                     ,1])                    
            isEven = False if isEven else True
            isNew =  False
    
        # loop on MH_ascendants

        max_n_level = 0     
        les_suivants = []
        temp_suivant = []
        
        for le_couple in MH_ascendants:
            n_level = le_couple.level - 1
            MH_adult1 = le_couple.adult1 
            if MH_adult1 != MH_none:

                max_n_level = max(max_n_level,n_level)

                MH_parents = []
                if n_level == n_level_max - 1: MH_parents = get_personne_parents(sql_obj,MH_adult1)

                la_display_ligne =[]

                la_display_ligne.append(col_format_levels[n_level][0] + [MH_adult1])
                la_display_ligne.append(col_format_levels[n_level][1] + [f'{n_level + n_rang}'])

                for k in range(n_level): 
                    la_display_ligne.append(col_format_levels[n_level][k+2] + ["⫷" if k == n_level-1 else " "])
                le_texte = f'<STRONG>{text_personne(MH_adult1,"prenom","nom")}</STRONG> {text_personne(MH_adult1,"prenoms","surnom")}'
                la_display_ligne.append(col_format_levels[n_level][2+n_level] + [le_texte])
                la_display_ligne.append(col_format_levels[n_level][3+n_level] + [f'{text_personne(MH_adult1,"byear")}'])
                la_display_ligne.append(col_format_levels[n_level][4+n_level] + [f'{text_personne(MH_adult1,"dyear")}'])
                la_display_ligne.append(col_format_levels[n_level][5+n_level] + [f'{text_personne(MH_adult1,"bcity")}'])
                la_display_ligne.append(col_format_levels[n_level][6+n_level] + [f'{text_personne(MH_adult1,"bregion")}'])
                la_display_ligne.append(col_format_levels[n_level][7+n_level] + [f'{text_personne(MH_adult1,"bcountry")}'])
                
                la_display_ligne.append(col_format_levels[n_level][8+n_level] + ["⫷" if MH_parents else " "])

                display_table.append([isNew,isEven,max_n_col,la_display_ligne,6])            
                isNew = False
                isEven = False if isEven else True

                if MH_parents : 
                    if MH_adult1.indi_id not in temp_suivant:
                        temp_suivant.append(MH_adult1.indi_id)
                        les_suivants.append(MH_adult1)

        if display_table:  
        #-------------------------------------------------------------------------------------------------------------          
            la_table = None
        #-------------------------------------------------------------------------------------------------------------
            for table_row in display_table:
        #-------------------------------------------------------------------------------------------------------------
                isNew = table_row[0]
        #---------------------------------------------------------------------------------------------------------- 
                if isNew:
                    la_page, boxes = PPTX_add_page(le_document,slide_layout_table_image) 
                    PPTX_add_box(la_page,0,0,ml,slide_height,bcolor = GraySide)

                    la_table = PPTX_table_add(la_page,col_format_level,ml,0,1,1)
                    if le_header : PPTX_table_row_display_photoID_in_cell(sql_obj,la_table,le_header,"header")  
        #----------------------------------------------------------------------------------------------------------    
                PPTX_table_row_display_photoID_in_cell(sql_obj,la_table,table_row)
            
        if les_suivants:
            
            for MH_item in les_suivants:
                PPTX_ascendant_table(sql_obj,le_document,MH_item,les_lignées,n_level_max,n_rang + max_n_level,15)

    return
#=============================================================================================================
import os
import psutil
import sys
sys.path.append(os.path.dirname(__file__))
from MH_MODULE_GEDCOM2SQL_1 import *
from MH_MODULE_SQLQUERIES_1 import *
from MH_MODULE_EXCEL_1 import *
from MH_MODULE_PPTX_1 import *
#--------------------------------------------------------------
from datetime import date,datetime
from pptx import Presentation
#--------------------------------------------------------------
Gedcom_dir  = "/Users/bernardconti/Downloads/"
icon_dir = icloud+'/MesProgrammes/python_global/icons'
save_path = "/Users/bernardconti/LOCAL_TEMP/Documents"
# --------------------------------------------------------------
"""

icloud = "/Users/bernardconti/Library/Mobile Documents/com~apple~CloudDocs"
les_WIP = "/Users/bernardconti/LOCAL_TEMP/WIP/"
les_PDF_RV = "/Users/bernardconti/Library/Mobile Documents/com~apple~CloudDocs/Les éditions du 57/PDF_RectoVerso/"

export_path =  "/Users/bernardconti/LOCAL_TEMP/Photos"
save_path = "/Users/bernardconti/LOCAL_TEMP/Documents"
import_dir = "/Users/bernardconti/LOCAL_TEMP/ePhoto_Doc/Import/"
watermarkdir = "/Users/bernardconti/LOCAL_TEMP/watermark"
le_MH_Photos_Bio = icloud+'/MesProgrammes/MH_Photos_Bio/'
#--------------------------------------------------------------
retour_ligne = '\n'
Gray= "#EFEEEF"  
GrayRow= "#EFEEEF"
Black ="#000000"
White = "#FFFFFF"
couleur_homme ="#ffcc99"
couleur_femme="#ccccff"
font_texte ="arial narrow"
couleur_cible = "#30DD30"
Gray1= "#b7b1b1"#F2F2F2
GraySide = "#c5c5c5"
Gray3 ="#D9D6D6"
Gray4 = "#E1DEDEEF"
couleur_chemin = "#EC5800"
couleur_titre ="#EB742F"
Green = "#B4E1C0"
Gray2= "#fde9d9"
Silver = "#A9A6A6" 
cell11 = "#F8CEB8"
cell12 = "#D9D9D9"
cell21 = "#FDE8DD"
cell22 = "#F2F2F2"
#Bleu = "#375e94"
Bleu = "#272170"
Darkblue = "#166082"
Gris_clair = "#EFEFEF"
ligne = "#375e94"
descendant = "#000000"
red = "#FF0000"""


#===========================================================================================
# DEBUT
#===========================================================================================
print("Execution de ",__file__)
#===========================================================================================
# Connect to the SQLite database (or create it if it doesn't exist)
#===========================================================================================
database_file = sqlite_gedcom2sql(Gedcom_dir)
#print (database_file)

connection_obj = sqlite3.connect(database_file)
connection_obj.row_factory = sqlite3.Row
sql_obj = connection_obj.cursor()

#===========================================================================================
#MENU
#===========================================================================================
from PyQt6.QtWidgets import *
from PyQt6.QtGui import *
from PyQt6.QtCore import *

custom_font = QFont("Arial",14)
line_height = 30
line_spacing = 10
cell_spacing = 5
n_line = 7
bio_width = [400,120,120,120,120,120]
bio_layout_width = 0
for number in bio_width:
    bio_layout_width += number
bio_layout_width += (len(bio_width)- 1)*cell_spacing
bio_layout_height = n_line * line_height + (n_line-1)*line_spacing

cell_width = 250
#========================================================================================    
for proc in psutil.process_iter():
    if proc.name() == "Microsoft PowerPoint": proc.kill()
#========================================================================================
# recupere tous les noms
#========================================================================================
list_noms = []
MH_individus = get_personne_all(sql_obj)
for MH_individu in MH_individus: 
    temp_nom = f"{MH_individu.nom} {MH_individu.prenom}{', '+MH_individu.prenoms if MH_individu.prenoms else ""} {'"' +MH_individu.surnom+'"' if MH_individu.surnom  else ""} | {MH_individu.indi_id}"
    list_noms.append(temp_nom)

list_noms = sorted(list_noms, key=lambda col0: col0[0])
#========================================================================================
# initiate list des directore pour impression PDF R&V
list_dirs = glob.glob(les_WIP+"*/", recursive = False)
if not list_dirs : exit()
#========================================================================================   
#  MAIN WINDOW   
#========================================================================================
class LA_WINDOW(QMainWindow):
#========================================================================================
    def __init__(self):
        super().__init__()
        self.setupUi()
#----------------------------------------------------------------------------------------
    def central_accueil(self):
        self.menu_clear_layout(self.central_layout)
        for item in self.infos:
            info = QLabel(item)
            info.setFixedSize(bio_layout_width, line_height)
            info.setStyleSheet("background-color: lightgray;color: black;"
                                    "qproperty-alignment: AlignLeft;"
                                    "border-radius: 10px;"
                                    "qproperty-wordWrap: true;"
                                    "padding: 6px;"
                                    )
            info.setFont(custom_font)
            self.central_layout.addRow(info) 
        self.menu_add_infos_row("2025, Created by Bernard CONTI, Dourdan, France")
#===========================================================================================================
    def setupUi(self):

        self.setObjectName("Menu_principal")
        self.resize(400, 400)

# Create menu bar, toolbar and statusbar objects
        menubar = self.menuBar()
        menubar.setFont(custom_font)
        
        self.centralwidget = QWidget(self)
        self.setCentralWidget(self.centralwidget)
        self.statusbar = QStatusBar(self)

# Create actions
        finAction = self.menu_add_action(icon = "exit.png",text = "Au revoir et Merci", cmd = self.PPTX_fin)
        excelAction = self.menu_add_action(text = "Arbre Ascendants Excel", cmd = self.central_excel)
        livresAction = self.menu_add_action(text = "Biographie", cmd = self.central_livres)
        BienvenueAction = self.menu_add_action(icon = "home.png", text = "Home Sweet Home",cmd = self.central_accueil)

# add  toolbar
        toolbar = QToolBar('Main ToolBar', self)
        toolbar.setIconSize(QSize(24, 24))
        toolbar.setStyleSheet("background-color: "+ couleur_femme +";color: black;"
                    "font-name=Arial;"
                    "font-size: 16px;"
                    )

        toolbar.addAction(BienvenueAction)
        toolbar.addAction(livresAction)
        toolbar.addAction(excelAction)
        toolbar.addActions(menubar.actions())
        toolbar.addAction(finAction)
        
# add toolbar and statusbar to main window
        self.addToolBar(toolbar)
        self.setStatusBar(self.statusbar)

# central_layout, self
        self.central_layout = QFormLayout(self.centralwidget) 
        self.central_layout.setVerticalSpacing(line_spacing)
        self.central_layout.setHorizontalSpacing(cell_spacing) 
#init central on welcome
        
        self.infos = ["Bienvenue dans le MIFA Studio"]
        self.central_accueil()

#======================================================================================================================          
    def PPTX_save(self):
        la_page = self.le_document.slides.add_slide(self.le_document.slide_layouts[slide_layout_57])
        OUT_fichier  = f'{save_path}/{self.OUTFICHIER.text()}'
        self.le_document.save(OUT_fichier)

        for file in os.listdir(watermarkdir):
            try: os.remove(watermarkdir+"/"+file)
            except Exception as error: print(error)

        self.infos.append("Livres > "+ OUT_fichier +  " enregistré")
        self.central_accueil()
#==============================================================================================================            
    def PPTX_fin(self):
        print("Au revoir et merci")
        self.force_close = True
        self.close()
#==============================================================================================================
# DEF_menu_add_...
#============================================================================================================== 
    def click_action(self, selected, deselected):

        for index in selected.indexes():
            item = self.tree.itemFromIndex(index)
            if self.tree.itemFromIndex(index.parent()) : 
                if item.text(0) not in self.list_selection_albums : 
                    self.list_selection_albums.append(item.text(0))
                    self.menu_add_infos_row(f'{self.tree.itemFromIndex(index.parent()).text(0)} > {item.text(0)} ')
            else:
                for f in self.les_folders:
                    if item.text(0) == f[0]:
                        for le_album in f[1].album_info:
                            if le_album.title not in self.list_selection_albums : 
                                if le_album.title[0] != "P": self.list_selection_albums.append(le_album.title)
                        self.menu_add_infos_row(f'Dossier : {item.text(0)} ')

#================================================================================================            
    def livres_tab_bio(self):
        self.menu_clear_layout(self.livres_bio_layout)
#--------------------------------------------------------------------------------------------------------------                      
# Input from list noms          
        i_row = 0
        selection_labels =  [
                    "La fratrie",              
                    "La cousinade", 
                    "Belle famille",     
                    "Tous les ascendants",       
                    "Les célébrités",
                    "Les biographies",   
                    "Les biographies simples",    
                    "Les Photos",
                    "Mode Fratrie"
                                ]
        self.cible = self.menu_add_input_text_completor(list_noms,width = bio_width[0],alignment="left")  #,popup = self.add_menu_personne
        self.livres_bio_layout.addRow(self.menu_add_label("Nom de la personne"),self.cible)

        self.personnes_options = [False]*(len(selection_labels))

        for i_row,selection_label in enumerate(selection_labels):
            self.personnes_options[i_row] = self.menu_add_checkbox(selection_label,n_cell = 1)
            self.livres_bio_layout.addRow(self.personnes_options[i_row]) #selection_label

        self.PPTX_action = self.menu_add_pushbutton("Publication Personne",self.publication_MH_personne)
        self.livres_bio_layout.addRow(self.PPTX_action)
#==============================================================================================================
    def livres_tab_document(self):
        self.menu_clear_layout(self.livres_document_layout)          
#--------------------------------------------------------------------------------------------------------------
# recupere tous les documents
        temp_documents = []
        for item in os.listdir(f'{save_path}/'):
            if item.split(".")[-1].lower() == "pptx" and "$" not in item:
                temp_documents.append(item)
        temp_documents = sorted(temp_documents, key=lambda col: (col[0]) )
        list_documents = ["Nouveau"] + temp_documents
#--------------------------------------------------------------------------------------------------------------
# recupere tous les models
        self.PPTX_label = self.menu_add_label("Choisir un fichier existant ou Initialiser un Nouveau PPTX",n_cell = 3)
        self.PPTX_file = self.menu_add_input_text_combobox(list_documents)  
        self.PPTX_action = self.menu_add_pushbutton("Suite ...",self.PPTX_Document)
        self.livres_document_layout.addRow(self.PPTX_label)
        self.livres_document_layout.addRow(self.PPTX_file,self.PPTX_action)
#===========================================================================================================
    def central_livres(self):
        self.menu_clear_layout(self.central_layout)

        self.livres_document,self.livres_document_layout = self.menu_add_layout("form")

        self.livres_tab = QTabWidget() 
        self.livres_tab.addTab(self.livres_document, 'Document')
        self.central_layout.addRow(self.livres_tab)

        self.livres_tab_document()
#===========================================================================================================          
    def central_excel(self):
        self.menu_clear_layout(self.central_layout)
        self.central_layout.addRow(self.menu_add_label("Sélectionnez la personne ...:",n_cell = 3))
        self.cible_x = self.menu_add_input_text_completor(list_noms)  
        self.central_layout.addRow(self.cible_x, self.menu_add_pushbutton("Suite ...",self.menu_Document_Excel, n_cell = 1))
#===========================================================================================================

# MENU
#===========================================================================================================   
    def menu_add_layout(self,type):
#--------------------------------------------------------------------------------------------------------------
        if type == "grid": item_layout = QGridLayout()
        else: item_layout = QFormLayout()

        item_layout.setVerticalSpacing(line_spacing)
        item_layout.setHorizontalSpacing(cell_spacing)  
        item = QWidget(self)
        item.setLayout(item_layout)
        return item,item_layout
#============================================================================================================== 
    def menu_clear_layout(self,le_layout):
#--------------------------------------------------------------------------------------------------------------
        while le_layout.count() > 0:
            item = le_layout.itemAt(0)
            widget = item.widget()
            if widget is None:
                le_layout.removeItem(item)
            else:
                #print(widget)
                widget.deleteLater()
                le_layout.removeWidget(widget)
        #print(le_layout.count())
        return
#==============================================================================================================         
    def menu_add_input_text_simple(self,default_input_text,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        if not default_input_text : return
        la_width = None
        n_cell = 2
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
        if not la_width : la_width = n_cell * cell_width
#-------------------------------------------------------------------------------------------------------------- 
        item = QLineEdit(default_input_text)     
        item.setStyleSheet("background-color: "+couleur_femme+";color: black")
        item.setAlignment(Qt.AlignmentFlag.AlignHCenter)
        item.setFont(custom_font) 
        item.setFixedSize(la_width,line_height) 
        return item
#==============================================================================================================
    def menu_add_input_text_completor(self,default_input_list,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        if not default_input_list : return
        la_width = None
        le_alignment = "center"
        n_cell = 2
        popup_action = None
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "popup"  : popup_action = valeur
                if clef == "alignment": le_alignment = valeur
        if not la_width : la_width = n_cell * cell_width

#--------------------------------------------------------------------------------------------------------------    
        le_completor = QCompleter(default_input_list)
        le_completor.setCaseSensitivity(Qt.CaseSensitivity.CaseInsensitive)
        if popup_action : le_completor.popup().clicked.connect(popup_action)
        le_completor.popup().setStyleSheet(
            "QListView {background-color: "+couleur_femme+";color: black;"
                        "selection-background-color: "+couleur_homme+";}"
            )
        item = QLineEdit(self)            
        item.setCompleter(le_completor)
        item.setStyleSheet("background-color: "+couleur_femme+";color: black")
        
        if le_alignment == "left":
            item.setAlignment(Qt.AlignmentFlag.AlignLeft)
        elif le_alignment == "right":
            item.setAlignment(Qt.AlignmentFlag.AlignRight)
        else:
            item.setAlignment(Qt.AlignmentFlag.AlignHCenter)
        item.setFont(custom_font) 
        item.setFixedSize(la_width,line_height) 
        
        return item     
#==============================================================================================================
    def menu_add_input_text_combobox(self,default_input_list,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        if not default_input_list : return
        la_width = None
        n_cell = 2
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item = QComboBox()
        item.addItems(default_input_list)      
        item.setStyleSheet("background-color: "+couleur_homme+";color: black")
        item.setFont(custom_font) 
        item.setFixedSize(la_width,line_height) 

        return item
#==============================================================================================================
    def menu_add_infos_row(self,le_texte):
#--------------------------------------------------------------------------------------------------------------
        self.statusbar.showMessage(le_texte)
        self.statusbar.setStyleSheet("background-color: "+couleur_homme+";color: black")
#==============================================================================================================
    def menu_add_label(self,le_label,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        la_width = None
        bcolor = "lightgray"
        color = "black"
        n_cell = 1
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "bcolor": bcolor = valeur
                if clef == "color": color = valeur

        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------                   
        item = QLabel(le_label)
        item.setStyleSheet("background-color: "+ bcolor +";color: "+color+";"
                                "qproperty-alignment: AlignLeft;"
                                "border-radius: 10px;"
                                "qproperty-wordWrap: true;"
                                "padding: 6px;"
                                )
        item.setFont(custom_font) 
        #item_label.setAlignment(Qt.AlignmentFlag.AlignHCenter)
        item.setFixedSize(la_width,line_height)

        return item
#==============================================================================================================
    def menu_add_date(self,**kwargs):
        la_width = None
        n_cell = 1
        format_osxphotos = "yyyy-MM-dd"
        format_edit = "dd-MM-yyyy"
        default = False

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                #if clef == "format":  format = valeur
                if clef == "default" : default = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item =QDateEdit()
        item.setDisplayFormat(format_edit)
        if default :
            item.setDate(QDate.fromString(default,format_edit))
        else :
            item.setDate(QDate.currentDate())
            
        item.setFont(custom_font)
        item.setFixedSize(la_width, line_height)
        item.setAlignment(Qt.AlignmentFlag.AlignHCenter)

        return item
#==============================================================================================================
    def menu_add_heure(self,**kwargs):
        la_width = None
        n_cell = 1
        default = False

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "default" : default = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item =QTimeEdit()
        item.setDisplayFormat('HH:mm')
        if default :
            item.setTime(QTime.fromString(default,'HH:mm:ss'))
        else :
            #item.setTime(QTime.currentTime())
            item.setTime(QTime(10,00,00))
            
        item.setFont(custom_font)
        item.setFixedSize(la_width, line_height)
        item.setAlignment(Qt.AlignmentFlag.AlignHCenter)

        return item
#==============================================================================================================
    def menu_add_pushbutton(self,le_texte,cmd,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        la_width = None
        n_cell = 1
        bcolor = "green"
        color = "white"

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "color" : color = valeur
                if clef == "bcolor" : bcolor = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item = QPushButton(le_texte)
        item.setStyleSheet(
                "QPushButton{background-color : "+bcolor+";color: "+color+"}"
                "QPushButton::pressed{background-color : red ;color: white}"
                            )
        item.setFont(custom_font)
        item.clicked.connect(cmd)
        item.setFixedSize(la_width, line_height)

        return item
#==============================================================================================================
    def menu_add_checkbox(self,le_texte,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        la_width = None
        n_cell = 1
        color_check = couleur_homme
        color = couleur_femme

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "color" : color = valeur
                if clef == "bcolor" : bcolor = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item = QCheckBox(le_texte)
        item.setStyleSheet(
            "QCheckBox{background-color : "+color+";color:black;}"
            "QCheckBox::checked{background-color : "+color_check+";color:black;}" )

        item.setFont(custom_font)
        item.setFixedSize(la_width, line_height)

        return item
#==============================================================================================================   
    def menu_add_action(self,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        tip = None
        cmd = None
        icon = None
        text = None

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "tip": tip = valeur
                if clef == "cmd": cmd = valeur
                if clef == "icon": icon = valeur
                if clef == "text": text = valeur
#--------------------------------------------------------------------------------------------------------------

        if icon and text :       item = QAction(QIcon(icon_dir+"/"+icon), text, self)
        elif not icon and text : item = QAction(text, self)
        elif icon and not text : item = QAction(QIcon(icon_dir+"/"+icon), "-",self)
        else: return None

        if cmd : item.triggered.connect(cmd)
        if tip : 
            #item.setToolTip(tip)
            item.setStatusTip(tip)
        item.setCheckable(True)
    
        return item
#==============================================================================================================           
                            
#==============================================================================================================
    def les_sub_folders(self,f_cur):
        temp_list = []

        for le_album in f_cur.album_info:

            if le_album.title[0] != "P": temp_list.append(f' > {le_album.title}')

        if temp_list   :   
            self.list_selection_albums.append(f_cur.title)
            self.list_folder_albums.append([f_cur.title,temp_list])
            self.list_selection_albums = self.list_selection_albums + temp_list

        for sf in f_cur.subfolders:
            self.les_sub_folders(sf)
        return 
#==============================================================================================================           
    def PPTX_reset(self):

        for item in self.personnes_options:
            for iitem in item:
                if iitem:
                    try: iitem.deleteLater() 
                    except : next 
        try:
            self.publication.deleteLater()
            self.mode_couples.deleteLater()
            self.mode_fratrie.deleteLater()
        except : next 
#==============================================================================================================                       
    def PPTX_Document(self):
#-------------------------------------------------------------------------------------------------------------- 
        le_fichier = self.PPTX_file.currentText()

        if le_fichier == "Nouveau":       
            ct=datetime.now()
            self.OUTFICHIER = f'LIVRE_{ct.year:04d}{ct.month:02d}{ct.day:02d}{ct.hour:02d}{ct.minute:02d}{ct.second:02d}.pptx'
# recupere tous les models

            list_models = []
            for item in os.listdir(f'{icloud}/MesProgrammes/Mes_Models/'):
                if item.split(".")[-1].lower() == "pptx" and item.startswith('MODEL_') :
                    list_models.append(item.replace(".pptx",""))
            list_models = sorted(list_models, key=lambda col: (col[0]) )
        
            if not list_models : 
                print("Pas de MODEL_")
                exit()
            
            self.livres_document_layout.addRow(self.menu_add_label("Model actif pour la création de Livre"))
            self.model_file = self.menu_add_input_text_combobox(list_models)
            self.livres_document_layout.addRow(self.model_file,self.menu_add_pushbutton("Suite ...",self.PPTX_add_new_document))
                                            
        else:

            IN_fichier  = f'{save_path}/{le_fichier}'
            self.OUTFICHIER = le_fichier
            print(IN_fichier)
            try :
                self.le_document = Presentation(IN_fichier)
                os.remove(IN_fichier)
            except Exception as error: 
                print(error)
                exit()

            self.menu_clear_layout(self.livres_document_layout)
            self.PPTX_label2 = self.menu_add_label("Vous pouvez modifier le nom du livre",n_cell = 3)
            self.livres_document_layout.addRow(self.PPTX_label2)
            self.OUTFICHIER = self.menu_add_input_text_simple(self.OUTFICHIER)
            self.PPTX_action = self.menu_add_pushbutton("Enregistrer le livre",self.PPTX_save)
            self.livres_document_layout.addRow(self.OUTFICHIER,self.PPTX_action)

            self.livres_bio,self.livres_bio_layout = self.menu_add_layout("grid")
            self.livres_tab.addTab(self.livres_bio,"Biographie")
            self.livres_tab_bio()

            self.livres_tab.setCurrentIndex(1)
#==============================================================================================================   
    def PPTX_add_new_document(self):
        
            self.le_document = Presentation(f'{icloud}/MesProgrammes/Mes_Models/{self.model_file.currentText()}.pptx') 

            self.menu_clear_layout(self.livres_document_layout)
            self.PPTX_label2 = self.menu_add_label("Vous pouvez modifier le nom du livre",n_cell = 3)
            self.livres_document_layout.addRow(self.PPTX_label2)
            self.OUTFICHIER = self.menu_add_input_text_simple(self.OUTFICHIER)
            self.PPTX_action = self.menu_add_pushbutton("Enregistrer le livre",self.PPTX_save)
            self.livres_document_layout.addRow(self.OUTFICHIER,self.PPTX_action)

            self.livres_bio,self.livres_bio_layout = self.menu_add_layout("form")
            self.livres_tab.addTab(self.livres_bio,"Biographie")
            self.livres_tab_bio()  

            self.livres_tab.setCurrentIndex(1)
#============================================================================================================== 
# S_EXCEL      
    def menu_Document_Excel(self):
#-------------------------------------------------------------------------------------------------------------- 
# label      
        self.excel_label = self.menu_add_label("Modifier le nom du fichier excel ...",n_cell = 3)     
        self.central_layout.addRow(self.excel_label)
# input & action
        la_cible = self.cible_x.text().split(" | ")
        la_cible_fullname = la_cible[0].replace(", ","_")
        la_cible_fullname = la_cible_fullname.replace(" ","_")
        la_cible_fullname = la_cible_fullname.replace('"',"")
        
        # type de tri
        self.n_level  = self.menu_add_input_text_simple("99",n_cell = 1)
        self.central_layout.addRow(self.menu_add_label("Nombre de niveau",n_cell = 2),self.n_level)

        self.output = self.menu_add_input_text_simple(f"ARBRE_{la_cible_fullname}.xlsx")
        self.excel_cmd = self.menu_add_pushbutton("Générer et Enregister",self.excel_gerenate_arbre)
        self.central_layout.addRow(self.output,self.excel_cmd)
#========================================================================================================
    def excel_gerenate_arbre(self):
#-------------------------------------------------------------------------------------------------------- 
        # Création de la sheet excel
        wb = Workbook() 
        ws_arbre = wb['Sheet'] 
        ws_arbre.title = 'Arbre'
        #---------------------------------------------------------           
        # génération de l'arbre  (MH_MODULE_EXCEL_1)
        la_cible = self.cible_x.text().split(" | ")
        la_cible_id = int(la_cible[1])
        n_levels = int(self.n_level.text())
        excel_arbre(sql_obj,ws_arbre,la_cible_id,n_levels) #MH_MODULE_EXCEL_1
        #---------------------------------------------------------
        # Enregistrement de l'arbre EXCEL et fin
        fname = f'{save_path}/{self.output.text()}'
        if os.path.isfile(fname): os.remove(fname)
        for proc in psutil.process_iter():
            if proc.name() == "Microsoft Excel": proc.kill()
        wb.save(fname)
        #---------------------------------------------------------
        # clear tab et back to Welcome
        self.infos.append("ARBRE > "+fname)
        self.central_accueil()
#=========================================================================================================
    def PPTX_Document_New(self):
#------------------------------------------------------------------------------------------------------------- 
        self.le_document = Presentation(f'{icloud}/MesProgrammes/Mes_Models/{self.model_file.currentText()}.pptx')

        self.PPTX_label.deleteLater()
        self.PPTX_file.deleteLater()
        self.PPTX_action.deleteLater()

        self.PPTX_label2 = self.menu_add_label("Modify PPTX File name",n_cell = 3)
        self.livres_document_layout.addRow(self.PPTX_label2)
        self.OUTFICHIER = self.menu_add_input_text_simple(self.OUTFICHIER)
        self.PPTX_action = self.menu_add_pushbutton("Save PPTX File",self.PPTX_save)
        self.livres_document_layout.addRow(self.OUTFICHIER,self.PPTX_action)

        self.livres_tab.addTab(self.livres_bio, 'Biographies')
        self.livres_tab.setCurrentIndex(1)  
#================================================================================================================= 
# S_INDI       
    def publication_MH_personne(self):
        #-----------------------------------------------------------------------------------------------------
        # Récupération de la cible
        la_cible = self.cible.text()
        MH_personne = get_personne_by_indi_id(sql_obj,int(la_cible.split(" | ")[1]))
        #-----------------------------------------------------------------------------------------------------
        # les options selectionnées                 
        isFratie = self.personnes_options[0].isChecked()
        isCouz = self.personnes_options[1].isChecked()
        isBellemifa = self.personnes_options[2].isChecked()
        isAscendants = self.personnes_options[3].isChecked()
        isCelebrity = self.personnes_options[4].isChecked()
        isBiographie = self.personnes_options[5].isChecked()
        isBiographieLight = self.personnes_options[6].isChecked()
        isPhotographies = self.personnes_options[7].isChecked()

        isModeFratrie = self.personnes_options[8].isChecked()

        temp_bro = get_personne_sisbros(sql_obj,MH_personne)
        if not temp_bro and isModeFratrie : isModeFratrie = False

        # les_bros inclus la personne et sa fratrie
        les_bros= [MH_personne]
        if isModeFratrie : les_bros= les_bros + temp_bro

        # Page de garde
        PPTX_add_page_garde(sql_obj,self.le_document,les_bros,isModeFratrie)
        
        #-------------------------------------------------------------------------------------------------------------           
        self.les_infos = ""       
        #----------------------------------------------------------------------------------------------------------------
        # Page sommaire
        la_page,boxes_sommaire = PPTX_page_image(sql_obj,self.le_document,MH_personne,
                                    "Sommaire","image_right",layout = slide_layout_section)
        box_sommaire = boxes_sommaire[1]
#------------------------------------------------------------------------------------------------------------
# PZOB
        n_level = 3
        MH_personnes = [] 

        PPTX_add_page_section(self.le_document,MH_personne,MH_none,"Parents, Grands-Parents et Arrières-Grands-Parents",box_sommaire)
        # Parents et Grands-parents
        titre=["Parents","Grands-Parents","Arrières-Grands-Parents"]   
        MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_personne,"ascendant","ascendant",n_level_max = n_level)  
        PPTX_ascendant_draw(sql_obj,self.le_document,MH_personne,n_level,titre,box_sommaire,MH_personnes,"single")

        for le_bro in les_bros:

            MH_conjoints = get_personne_conjoints(sql_obj,le_bro)

            #Biographies personnelles
            PPTX_add_page_section(self.le_document,le_bro,MH_none,"Biographie",box_sommaire)
            PPTX_biographies_table(sql_obj,self.le_document,[[le_bro,"cible",le_bro,0,""]],"main",None)
            if isPhotographies : PPTX_personne_MHphotos(sql_obj,self.le_document,le_bro)

            # Conjoints,Enfants,Petits-Enfants
            PPTX_add_page_section(self.le_document,le_bro,MH_none,"Le foyer",box_sommaire)
            titre=["Conjoints","Enfants","Petits-Enfants","Arrières-Petits-Enfants"] 
            MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,le_bro,"descendant","descendant",n_level_max = 99)   
            PPTX_descendant_group_table(sql_obj,self.le_document,le_bro,titre,99,None,"descendant") 
            PPTX_descendant_arbre_table(sql_obj,self.le_document,le_bro,[MH_none,MH_none],titre,99,None,"descendant",20)           
                                                                                    
        if isFratie and not isModeFratrie:
            # Frères ou Soeurs,Neuveux,Petits-Neuveux
            titre=["Frères ou Soeurs","Neuveux","Petits-Neuveux"]
            MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_personne,"descendant","fratrie",n_level_max = n_level)
            PPTX_descendant_group_table(sql_obj,self.le_document,MH_personne,titre,n_level,box_sommaire,"fratrie") 
            PPTX_descendant_arbre_table(sql_obj,self.le_document,MH_personne,[MH_none,MH_none],titre,99,None,"fratrie",20)

        if isCouz:
            # Oncles paternels,Cousins,Petits-Cousins
            titre=["Oncles paternels","Cousins","Petits-Cousins"]
            MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_personne,"descendant","couzpater",n_level_max = n_level)
            PPTX_descendant_group_table(sql_obj,self.le_document,MH_personne,titre,n_level,box_sommaire,"couzpater") 
            PPTX_descendant_arbre_table(sql_obj,self.le_document,MH_personne,[MH_none,MH_none],titre,99,None,"couzpater",20)

            # Oncles maternel,Cousins,Petits-Cousins
            titre=["Oncles maternels","Cousins","Petits-Cousins"]  
            MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_personne,"descendant","couzmater",n_level_max = n_level) 
            PPTX_descendant_group_table(sql_obj,self.le_document,MH_personne,titre,n_level,box_sommaire,"couzmater") 
            PPTX_descendant_arbre_table(sql_obj,self.le_document,MH_personne,[MH_none,MH_none],titre,99,None,"couzmater",20)

        if isBellemifa:

            for le_bro in les_bros:
                MH_conjoints = get_personne_conjoints(sql_obj,le_bro)
                if MH_conjoints :
                    # Beaux Parents et Grands-parents
                    for MH_conjoint in MH_conjoints:
                        PPTX_add_page_section(self.le_document,MH_conjoint,MH_none,"Entourage",box_sommaire)
                        MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_conjoint,"ascendant","ascendant",n_level_max = n_level)
                        PPTX_biographies_table(sql_obj,self.le_document,[[MH_conjoint,"cible",MH_conjoint,0,""]],"main",None)
                        titre=["Beaux parents","Grands-Parents","Arrières-Grands-Parents"]    
                        PPTX_ascendant_draw(sql_obj,self.le_document,MH_conjoint,n_level,titre,box_sommaire,MH_personnes,"single")  
                    # Beaux Frères ou Soeurs,Neuveux,Petits-Neuveux
                        titre=["Beaux Frères ou Soeurs","Neuveux","Petits-Neuveux"] 
                        MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_conjoint,"descendant","fratrie",n_level_max = n_level)   
                        PPTX_descendant_group_table(sql_obj,self.le_document,MH_conjoint,titre,n_level,box_sommaire,"fratrie")
                        PPTX_descendant_arbre_table(sql_obj,self.le_document,MH_conjoint,[MH_none,MH_none],titre,99,None,"fratrie",20)
            
        if isCelebrity : 
            MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_personne,"descendant","descendant") 
            MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_personne,"descendant","fratrie") 
            MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_personne,"descendant","couzpater") 
            MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_personne,"descendant","couzmater") 
            MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_personne,"ascendant","ascendant") 
            for le_bro in les_bros:
                MH_conjoints = get_personne_conjoints(sql_obj,le_bro)
                if MH_conjoints :
                    for MH_conjoint in MH_conjoints:
                        MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_conjoint,"ascendant","ascendant") 
                        MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_conjoint,"descendant","fratrie") 
                        MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_conjoint,"descendant","couzpater") 
                        MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_conjoint,"descendant","couzmater") 
            PPTX_add_page_section(self.le_document,MH_none,MH_none,"Célébrités",box_sommaire)
            PPTX_biographies_table(sql_obj,self.le_document,MH_personnes,"celebrity",box_sommaire)

        if isBiographie: 
            PPTX_add_page_section(self.le_document,MH_none,MH_none,"Biographies complètes",box_sommaire)
            PPTX_biographies_table(sql_obj,self.le_document,MH_personnes,"bio",None)

        if isBiographieLight: 
            PPTX_add_page_section(self.le_document,MH_none,MH_none,"Biographies simples",box_sommaire)
            PPTX_biographies_light_table(sql_obj,self.le_document,MH_personnes,None)

        if isAscendants: 
            self.les_infos = self.les_infos + " • Ascendants" 
            PPTX_add_page_section(self.le_document,MH_personne,MH_none,"Liste détaillée des ascendants",box_sommaire)
            MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_personne,"ascendant","ascendant") 
            MH_lignées = get_personne_lignées(sql_obj,MH_personne)
            PPTX_ascendant_table(sql_obj,self.le_document,MH_personne,MH_lignées,4,0,12.5)
            #for le_bro in les_bros:
            #    MH_conjoints = get_personne_conjoints(sql_obj,le_bro)
            #    if MH_conjoints:
            #        for MH_conjoint in MH_conjoints:
            #            PPTX_add_page_section(self.le_document,MH_conjoint,MH_none,"Liste détaillée des ascendants",box_sommaire)
            #            MH_personnes = PPTX_feed_MH_personnes(sql_obj,MH_personnes,MH_conjoint,"ascendant","ascendant") 
            #            MH_lignées = get_personne_lignées(sql_obj,MH_conjoint)
            #            PPTX_ascendant_table(sql_obj,self.le_document,MH_lignées,MH_conjoint,[MH_conjoint],0)

#--------------------------------------------------------------------------------------                           
# on nettoie la page
#--------------------------------------------------------------------------------------
        self.livres_tab.setCurrentIndex(0)
        self.menu_add_infos_row("Biobraphies > "+self.les_infos)
        self.livres_tab_bio()
#==============================================================================================================                           
#==============================================================================================================   
if __name__ == "__main__":
    import sys
    app = QApplication(sys.argv)
    MainWindow = LA_WINDOW()
    MainWindow.show()
    sys.exit(app.exec())
#==============================================================================================================   
#==============================================================================================================                           
#==============================================================================================================   
